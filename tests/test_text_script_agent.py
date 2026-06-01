from __future__ import annotations

import json
import os
import time
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any
from unittest.mock import patch

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

from ppt_system.export.export_layer_mode import SEPARATE_LAYER_MODE
from ppt_system.export.export_page_resume import CHECKPOINT_FILE_NAME
from ppt_system.export.export_step_checkpoint import STEP_CHECKPOINT_DIR_NAME
from ppt_system.export.direct_page_script import (
    build_direct_page_refine_prompt,
    build_direct_page_prompt,
    prepare_direct_page_assets,
)
from ppt_system.export.export_pipeline import export_project_to_pptx
from ppt_system.export.text_script_runtime import build_project_script_source, execute_generated_text_script, normalize_page_script
from ppt_system.export.text_style_runtime import should_wrap_text


class FakeChatProvider:
    def __init__(self, responses: list[dict[str, Any]]) -> None:
        self.responses = list(responses)
        self.calls: list[list[dict[str, Any]]] = []
        self.api_base_url = "https://example.com/v1"
        self.model = "fake-chat"
        self.temperature = 0.3
        self.max_tokens = 5000
        self.reasoning_effort = ""

    def build_image_message_item(self, image_path: Path) -> dict[str, Any]:
        return {"type": "image_url", "image_url": {"url": str(image_path)}}

    def complete_json(self, messages: list[dict[str, Any]]) -> dict[str, Any]:
        self.calls.append(messages)
        if not self.responses:
            raise RuntimeError("no more responses")
        return self.responses.pop(0)


def _write_minimal_assets_manifest(page_assets_dir: Path, *, image_width: int, image_height: int) -> None:
    page_assets_dir.mkdir(parents=True, exist_ok=True)
    Image.new("RGBA", (10, 10), (0, 128, 255, 255)).save(page_assets_dir / "asset_001.png")
    (page_assets_dir / "assets.json").write_text(
        json.dumps(
            {
                "image_width": image_width,
                "image_height": image_height,
                "assets": [
                    {"index": 1, "file": "asset_001.png", "left": 20, "top": 20, "width": 40, "height": 40}
                ],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )


class TextScriptRuntimeAndDirectPathTests(unittest.TestCase):
    def test_should_wrap_text_keeps_single_line_banner_and_badge_unwrapped(self) -> None:
        self.assertFalse(should_wrap_text("01", 118, 88, 38))
        self.assertFalse(should_wrap_text("AI 转换", 300, 58, 26))
        self.assertFalse(should_wrap_text("进阶玩法：让 AI 从聊天助手变成流程生产工具", 1350, 82, 30))

    def test_should_wrap_text_allows_explicit_multiline_and_tall_body_wrapping(self) -> None:
        self.assertTrue(should_wrap_text("第一行\n第二行", 240, 90, 20))
        self.assertTrue(should_wrap_text("这是一个需要自动换行的长段落文本", 180, 140, 20))

    def test_direct_page_prompt_keeps_reference_and_elements_constraints(self) -> None:
        prompt = build_direct_page_prompt(
            image_width=2048,
            image_height=1152,
            text_placeholders={
                "placeholders": [
                    {
                        "id": "text_01",
                        "left": 100,
                        "top": 80,
                        "width": 420,
                        "height": 70,
                        "font_size": 28,
                        "color": "123A63",
                        "align": "LEFT",
                    }
                ]
            },
        )
        self.assertIn("第一张图是完整原稿图，第二张图是去文字后的元素图", prompt)
        self.assertIn("text_placeholders", prompt)
        self.assertIn("默认沿用 placeholder", prompt)
        self.assertIn("元素会在导出时单独加入", prompt)

    def test_direct_page_refine_prompt_allows_asset_adjustments(self) -> None:
        prompt = build_direct_page_refine_prompt(
            image_width=2048,
            image_height=1152,
            page_script='add_text(slide, "标题", 100, 100, 300, 60, size=24)',
            asset_adjustments={"global": {"dy": 6}},
            round_index=0,
        )
        self.assertIn("asset_adjustments", prompt)
        self.assertIn("本轮只修文字", prompt)
        self.assertIn("asset_adjustments 固定返回空对象 {}", prompt)
        self.assertIn("请直接修正 page_script", prompt)

    def test_prepare_assets_exposes_global_alignment_as_default_asset_adjustment(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            reference_path = root / "reference.png"
            visual_path = root / "visual.png"

            reference = Image.new("RGBA", (240, 140), (255, 255, 255, 255))
            draw_reference = ImageDraw.Draw(reference)
            draw_reference.rectangle((80, 40, 140, 90), outline=(0, 80, 220, 255), width=4)
            reference.save(reference_path)

            visual = Image.new("RGBA", (240, 140), (255, 255, 255, 255))
            draw_visual = ImageDraw.Draw(visual)
            draw_visual.rectangle((55, 58, 115, 108), outline=(0, 80, 220, 255), width=4)
            visual.save(visual_path)

            result = prepare_direct_page_assets(
                work_dir=work_dir,
                page_no=1,
                elements_image=visual_path,
                reference_image=reference_path,
                reference_text_boxes=[],
                image_width=240,
                image_height=140,
            )

            self.assertIsNotNone(result.global_alignment)
            self.assertTrue(bool(result.global_alignment["should_apply"]))
            self.assertEqual(
                result.asset_adjustments,
                {
                    "global": {
                        "dx": int(result.global_alignment["dx"]),
                        "dy": int(result.global_alignment["dy"]),
                    }
                },
            )

    def test_prepare_direct_page_assets_preserves_existing_transparent_input_and_tiny_components(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            visual_path = root / "visual.png"

            image = Image.new("RGBA", (40, 40), (255, 255, 255, 0))
            draw = ImageDraw.Draw(image)
            draw.rectangle((5, 5, 15, 15), fill=(0, 80, 220, 255))
            image.putpixel((30, 30), (0, 80, 220, 255))
            image.save(visual_path)

            result = prepare_direct_page_assets(
                work_dir=work_dir,
                page_no=1,
                elements_image=visual_path,
                reference_image=None,
                image_width=40,
                image_height=40,
            )

            self.assertEqual(result.split_source_image, str(visual_path))
            self.assertEqual(result.transparent_preview_image, str(visual_path))
            self.assertEqual(int(result.manifest["min_area"]), 1)
            self.assertEqual(int(result.manifest["count"]), 2)
            self.assertFalse((work_dir / "page_01" / "page_01_enhanced.png").exists())
            self.assertFalse((work_dir / "page_01" / "page_01_transparent.png").exists())

    def test_normalize_page_script_rejects_disallowed_code(self) -> None:
        with self.assertRaises(RuntimeError):
            normalize_page_script("import os")

    def test_normalize_page_script_preserves_valid_text_call(self) -> None:
        script = 'add_text(slide, "标题", 10, 20, 200, 60, size=24, color="163A63")'
        self.assertEqual(normalize_page_script(script), script)

    def test_normalize_page_script_coalesces_multiline_call(self) -> None:
        script = 'add_text(slide,\n"标题",\n10, 20, 200, 60,\nsize=24, color="163A63")'
        expected = 'add_text(slide, "标题", 10, 20, 200, 60, size=24, color="163A63")'
        self.assertEqual(normalize_page_script(script), expected)

    def test_normalize_page_script_accepts_json_style_literals_in_add_runs(self) -> None:
        script = (
            'add_runs(slide, [{"text":"从","size":64,"color":"08265C","bold":true},'
            '{"text":"提问","size":64,"color":"0B55E6","bold":false},'
            '{"text":"到","size":64,"color":"08265C","bold":true},'
            '{"text":"协同","size":64,"color":"0B55E6","italic":null}],'
            '100, 120, 500, 90, align="LEFT", anchor="TOP")'
        )
        expected = (
            'add_runs(slide, [{"text": "从", "size": 64, "color": "08265C", "bold": True}, '
            '{"text": "提问", "size": 64, "color": "0B55E6", "bold": False}, '
            '{"text": "到", "size": 64, "color": "08265C", "bold": True}, '
            '{"text": "协同", "size": 64, "color": "0B55E6", "italic": None}], '
            '100, 120, 500, 90, align="LEFT", anchor="TOP")'
        )
        self.assertEqual(normalize_page_script(script), expected)

    def test_build_and_execute_project_script_produces_editable_ppt(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            page_assets_dir = work_dir / "page_01" / "assets"
            page_assets_dir.mkdir(parents=True, exist_ok=True)
            output_pptx = root / "result.pptx"
            Image.new("RGBA", (40, 30), (0, 128, 255, 255)).save(page_assets_dir / "asset_001.png")
            (page_assets_dir / "assets.json").write_text(
                json.dumps(
                    {"assets": [{"file": "asset_001.png", "left": 40, "top": 50, "width": 120, "height": 80}]},
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            project = {
                "slide_width_inch": 13.333333,
                "image_width": 400,
                "image_height": 240,
                "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
                "pages": [
                    {
                        "page_no": 1,
                        "title": "示例页",
                        "summary": "摘要",
                        "texts": [
                            {"role": "title", "text": "示例页", "left": 20, "top": 20, "width": 160, "height": 40},
                            {"role": "body", "text": "第一条", "left": 30, "top": 120, "width": 200, "height": 40},
                        ],
                    }
                ],
            }
            script_source = build_project_script_source(
                project,
                work_dir,
                output_pptx,
                [{"page_no": 1, "script": 'add_text_ref(slide, page_texts, "title", 18, 18, 180, 44, size=24, color="163A63", bold=True)\nadd_text_ref(slide, page_texts, "body_1", 36, 122, 210, 36, size=14, color="355C7D", bold=False)'}],
                include_assets=True,
            )
            script_path = work_dir / "generated_text_layout.py"
            script_path.write_text(script_source, encoding="utf-8")

            execute_generated_text_script(script_path)

            prs = Presentation(str(output_pptx))
            slide = prs.slides[0]
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            self.assertIn("示例页", texts)
            self.assertIn("第一条", texts)
            text_shapes = [shape for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            self.assertEqual(text_shapes[0].text_frame.auto_size, MSO_AUTO_SIZE.NONE)

    def test_execute_generated_script_resolves_project_imports_from_external_cwd(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            script_dir = root / "page_01"
            script_dir.mkdir(parents=True, exist_ok=True)
            isolated_cwd = root / "external_cwd"
            isolated_cwd.mkdir()
            output_path = script_dir / "worker_result.txt"
            script_path = script_dir / "generated_text_layout_preview_round_01.py"
            script_path.write_text(
                """
from __future__ import annotations

from pathlib import Path

from ppt_system.export.text_style_runtime import should_wrap_text


def build_deck():
    output_path = Path(__file__).with_name("worker_result.txt")
    output_path.write_text(str(should_wrap_text("正文内容", 80, 120, 18)), encoding="utf-8")
    return output_path
""".lstrip(),
                encoding="utf-8",
            )

            original_cwd = Path.cwd()
            try:
                os.chdir(isolated_cwd)
                execute_generated_text_script(script_path)
            finally:
                os.chdir(original_cwd)

            self.assertTrue(output_path.exists())
            self.assertIn(output_path.read_text(encoding="utf-8"), {"True", "False"})

    def test_execute_generated_script_stops_running_worker(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            marker_path = root / "started.txt"
            output_path = root / "worker_result.txt"
            script_path = root / "slow_generated_text_layout.py"
            script_path.write_text(
                f"""
from __future__ import annotations

import time
from pathlib import Path


def build_deck():
    Path(r"{marker_path}").write_text("started", encoding="utf-8")
    time.sleep(30)
    Path(r"{output_path}").write_text("done", encoding="utf-8")
    return Path(r"{output_path}")
""".lstrip(),
                encoding="utf-8",
            )
            checks = {"count": 0}

            def stop_after_worker_starts() -> bool:
                checks["count"] += 1
                return marker_path.exists()

            with self.assertRaises(InterruptedError):
                execute_generated_text_script(
                    script_path,
                    timeout_seconds=10,
                    stop_checker=stop_after_worker_starts,
                )

            self.assertFalse(output_path.exists())
            self.assertGreater(checks["count"], 0)

    def test_execute_generated_script_stops_child_process_tree(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            marker_path = root / "started.txt"
            child_marker_path = root / "child_started.txt"
            output_path = root / "child_done.txt"
            child_script_path = root / "slow_child.py"
            script_path = root / "slow_parent_generated_text_layout.py"
            child_script_path.write_text(
                f"""
from __future__ import annotations

import time
from pathlib import Path


Path(r"{child_marker_path}").write_text("started", encoding="utf-8")
time.sleep(30)
Path(r"{output_path}").write_text("done", encoding="utf-8")
""".lstrip(),
                encoding="utf-8",
            )
            script_path.write_text(
                f"""
from __future__ import annotations

import subprocess
import sys
import time
from pathlib import Path


def build_deck():
    Path(r"{marker_path}").write_text("started", encoding="utf-8")
    subprocess.Popen([sys.executable, r"{child_script_path}"])
    time.sleep(30)
    return Path(r"{output_path}")
""".lstrip(),
                encoding="utf-8",
            )

            def stop_after_child_starts() -> bool:
                return child_marker_path.exists()

            with self.assertRaises(InterruptedError):
                execute_generated_text_script(
                    script_path,
                    timeout_seconds=10,
                    stop_checker=stop_after_child_starts,
                )

            time.sleep(0.5)
            self.assertFalse(output_path.exists())

    def test_generated_script_can_split_assets_and_texts_into_two_slides(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            page_assets_dir = work_dir / "page_01" / "assets"
            page_assets_dir.mkdir(parents=True, exist_ok=True)
            output_pptx = root / "result.pptx"
            Image.new("RGBA", (40, 30), (0, 128, 255, 255)).save(page_assets_dir / "asset_001.png")
            (page_assets_dir / "assets.json").write_text(
                json.dumps(
                    {
                        "image_width": 400,
                        "image_height": 240,
                        "assets": [
                            {"index": 1, "file": "asset_001.png", "left": 40, "top": 50, "width": 120, "height": 80}
                        ],
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            project = {
                "slide_width_inch": 10.0,
                "image_width": 400,
                "image_height": 240,
                "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
                "pages": [
                    {
                        "page_no": 1,
                        "title": "示例页",
                        "summary": "摘要",
                        "texts": [{"role": "title", "text": "示例页", "left": 40, "top": 50, "width": 120, "height": 80}],
                    }
                ],
            }
            script_source = build_project_script_source(
                project,
                work_dir,
                output_pptx,
                [{"page_no": 1, "script": 'add_text_ref(slide, page_texts, "title", 40, 50, 120, 80, size=24, color="163A63", bold=True)'}],
                include_assets=True,
                layer_mode=SEPARATE_LAYER_MODE,
            )
            script_path = work_dir / "generated_text_layout.py"
            script_path.write_text(script_source, encoding="utf-8")

            execute_generated_text_script(script_path)

            prs = Presentation(str(output_pptx))
            self.assertEqual(len(prs.slides), 2)
            asset_slide = prs.slides[0]
            text_slide = prs.slides[1]

            asset_picture = [shape for shape in asset_slide.shapes if shape.shape_type == 13][0]
            self.assertEqual(len([shape for shape in asset_slide.shapes if hasattr(shape, "text") and shape.text]), 0)

            text_box = [shape for shape in text_slide.shapes if hasattr(shape, "text") and shape.text][0]
            self.assertEqual(len([shape for shape in text_slide.shapes if shape.shape_type == 13]), 0)
            self.assertEqual(text_box.text, "示例页")
            self.assertAlmostEqual(asset_picture.left / prs.slide_width, text_box.left / prs.slide_width, places=3)
            self.assertAlmostEqual(asset_picture.top / prs.slide_height, text_box.top / prs.slide_height, places=3)
            self.assertAlmostEqual(asset_picture.width / prs.slide_width, text_box.width / prs.slide_width, places=3)
            self.assertAlmostEqual(asset_picture.height / prs.slide_height, text_box.height / prs.slide_height, places=3)

    def test_generated_script_maps_assets_using_manifest_canvas_size(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            page_assets_dir = work_dir / "page_01" / "assets"
            page_assets_dir.mkdir(parents=True, exist_ok=True)
            output_pptx = root / "result.pptx"
            Image.new("RGBA", (10, 10), (0, 128, 255, 255)).save(page_assets_dir / "asset_001.png")
            (page_assets_dir / "assets.json").write_text(
                json.dumps(
                    {
                        "image_width": 800,
                        "image_height": 400,
                        "assets": [
                            {"file": "asset_001.png", "left": 80, "top": 40, "width": 160, "height": 80}
                        ],
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            project = {
                "slide_width_inch": 10.0,
                "image_width": 400,
                "image_height": 200,
                "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
                "pages": [{"page_no": 1, "title": "示例页", "summary": "", "texts": []}],
            }
            script_source = build_project_script_source(
                project,
                work_dir,
                output_pptx,
                [{"page_no": 1, "script": ""}],
                include_assets=True,
            )
            script_path = work_dir / "generated_text_layout.py"
            script_path.write_text(script_source, encoding="utf-8")

            execute_generated_text_script(script_path)

            prs = Presentation(str(output_pptx))
            picture_shapes = [shape for shape in prs.slides[0].shapes if shape.shape_type == 13]
            self.assertEqual(len(picture_shapes), 1)
            picture = picture_shapes[0]
            self.assertAlmostEqual(picture.left / prs.slide_width, 0.1, places=3)
            self.assertAlmostEqual(picture.top / prs.slide_height, 0.1, places=3)

    def test_generated_script_applies_asset_adjustments(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            page_assets_dir = work_dir / "page_01" / "assets"
            page_assets_dir.mkdir(parents=True, exist_ok=True)
            output_pptx = root / "result.pptx"
            Image.new("RGBA", (10, 10), (0, 128, 255, 255)).save(page_assets_dir / "asset_001.png")
            (page_assets_dir / "assets.json").write_text(
                json.dumps(
                    {
                        "image_width": 400,
                        "image_height": 200,
                        "assets": [
                            {"index": 1, "file": "asset_001.png", "left": 40, "top": 20, "width": 80, "height": 40}
                        ],
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            project = {
                "slide_width_inch": 10.0,
                "image_width": 400,
                "image_height": 200,
                "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
                "pages": [{"page_no": 1, "title": "示例页", "summary": "", "texts": []}],
            }
            script_source = build_project_script_source(
                project,
                work_dir,
                output_pptx,
                [{"page_no": 1, "script": "", "asset_adjustments": {"global": {"dx": 20, "dy": 10}}}],
                include_assets=True,
            )
            script_path = work_dir / "generated_text_layout.py"
            script_path.write_text(script_source, encoding="utf-8")

            execute_generated_text_script(script_path)

            prs = Presentation(str(output_pptx))
            picture = [shape for shape in prs.slides[0].shapes if shape.shape_type == 13][0]
            self.assertAlmostEqual(picture.left / prs.slide_width, 0.15, places=3)
            self.assertAlmostEqual(picture.top / prs.slide_height, 0.15, places=3)

    def test_export_project_to_pptx_uses_direct_office_refine(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            output_pptx = root / "result.pptx"
            visual_path = root / "visual.png"
            reference_path = root / "reference.png"
            Image.new("RGBA", (400, 240), (255, 255, 255, 0)).save(visual_path)
            with Image.open(visual_path).convert("RGBA") as image:
                image.paste((0, 82, 214, 255), (40, 40, 120, 100))
                image.save(visual_path)
            Image.new("RGBA", (400, 240), (255, 255, 255, 255)).save(reference_path)
            project = {
                "slide_width_inch": 13.333333,
                "image_width": 400,
                "image_height": 240,
                "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
                "pages": [
                    {
                        "page_no": 1,
                        "title": "脚本页",
                        "summary": "摘要",
                        "visual_image": str(visual_path),
                        "reference_image": str(reference_path),
                        "texts": [],
                    }
                ],
            }
            provider = FakeChatProvider([{"page_script": 'add_text(slide, "脚本页", 12, 14, 130, 36, size=20, color="163A63", bold=True)'}])

            with patch("ppt_system.export.direct_project_script.render_pptx_first_slide_to_png", return_value=None):
                result = export_project_to_pptx(
                    project,
                    work_dir,
                    output_pptx,
                    chat_provider=provider,  # type: ignore[arg-type]
                )

            self.assertEqual(result["text_layout_strategy"], "direct_office_refine")
            self.assertTrue(Path(result["text_script_path"]).exists())
            self.assertTrue(output_pptx.exists())
            self.assertIn("page_results", result)
            self.assertEqual(result["layer_mode"], SEPARATE_LAYER_MODE)
            self.assertEqual(result["delivery_mode"], "separate_layer_slides")
            self.assertEqual(result["logical_page_count"], 1)
            self.assertEqual(result["page_count"], 2)
            prs = Presentation(str(output_pptx))
            self.assertEqual(len(prs.slides), 2)

    def test_export_project_keeps_existing_asset_adjustments_after_overlap_check(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            output_pptx = root / "result.pptx"
            visual_path = root / "visual.png"
            reference_path = root / "reference.png"
            Image.new("RGBA", (400, 240), (255, 255, 255, 0)).save(visual_path)
            with Image.open(visual_path).convert("RGBA") as image:
                image.paste((0, 82, 214, 255), (40, 40, 120, 100))
                image.save(visual_path)
            Image.new("RGBA", (400, 240), (255, 255, 255, 255)).save(reference_path)
            project = {
                "slide_width_inch": 13.333333,
                "image_width": 400,
                "image_height": 240,
                "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
                "pages": [
                    {
                        "page_no": 1,
                        "title": "脚本页",
                        "summary": "摘要",
                        "visual_image": str(visual_path),
                        "reference_image": str(reference_path),
                        "texts": [],
                    }
                ],
            }
            provider = FakeChatProvider(
                [
                    {"page_script": 'add_text(slide, "首轮文字", 12, 14, 130, 36, size=20, color="163A63", bold=True)'},
                    {
                        "page_script": 'add_text(slide, "二轮改字", 20, 24, 150, 40, size=22, color="163A63", bold=True)',
                    },
                ]
            )
            _write_minimal_assets_manifest(work_dir / "page_01" / "assets", image_width=400, image_height=240)

            fake_asset_result = type(
                "FakePreparedAssets",
                (),
                {
                    "manifest_path": str(work_dir / "page_01" / "assets" / "assets.json"),
                    "manifest": {},
                    "image_width": 400,
                    "image_height": 240,
                    "split_source_image": str(visual_path),
                    "transparent_preview_image": str(visual_path),
                    "removed_intermediate_images": [],
                    "global_alignment": {"should_apply": True, "dx": -10, "dy": 40},
                    "asset_adjustments": {"global": {"dx": -10, "dy": 40}},
                },
            )()

            with patch("ppt_system.export.direct_project_script.render_pptx_first_slide_to_png", return_value=reference_path):
                with patch("ppt_system.export.direct_project_script.prepare_direct_page_assets", return_value=fake_asset_result):
                    with patch(
                        "ppt_system.export.direct_project_script.analyze_text_asset_overlaps",
                        return_value=type(
                            "FakeOverlap",
                            (),
                            {
                                "total_boxes": 1,
                                "overlap_box_count": 1,
                                "overlap_ratio": 1.0,
                                "max_overlap_pixels": 200,
                                "overlapping_box_indices": [1],
                            },
                        )(),
                    ):
                        result = export_project_to_pptx(
                            project,
                            work_dir,
                            output_pptx,
                            chat_provider=provider,  # type: ignore[arg-type]
                        )

            final_script = Path(result["text_script_path"]).read_text(encoding="utf-8")
            self.assertIn("二轮改字", final_script)
            self.assertIn('"dx": -10', final_script)
            self.assertIn('"dy": 40', final_script)
            self.assertNotIn("asset_alignment_decision", result["page_results"][0])

    def test_export_project_to_pptx_requires_chat_provider(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            output_pptx = root / "result.pptx"
            visual_path = root / "visual.png"
            reference_path = root / "reference.png"
            Image.new("RGBA", (400, 240), (255, 255, 255, 255)).save(visual_path)
            Image.new("RGBA", (400, 240), (255, 255, 255, 255)).save(reference_path)
            project = {
                "slide_width_inch": 13.333333,
                "image_width": 400,
                "image_height": 240,
                "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
                "pages": [
                    {
                        "page_no": 1,
                        "title": "脚本页",
                        "summary": "摘要",
                        "visual_image": str(visual_path),
                        "reference_image": str(reference_path),
                        "texts": [],
                    }
                ],
            }

            with self.assertRaises(RuntimeError):
                export_project_to_pptx(project, work_dir, output_pptx)

    def test_export_project_to_pptx_resumes_completed_pages_from_checkpoints(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            output_pptx = root / "result.pptx"
            visual_path_1 = root / "visual_01.png"
            reference_path_1 = root / "reference_01.png"
            visual_path_2 = root / "visual_02.png"
            reference_path_2 = root / "reference_02.png"
            for path in (visual_path_1, visual_path_2):
                Image.new("RGBA", (400, 240), (255, 255, 255, 0)).save(path)
                with Image.open(path).convert("RGBA") as image:
                    image.paste((0, 82, 214, 255), (40, 40, 120, 100))
                    image.save(path)
            for path in (reference_path_1, reference_path_2):
                Image.new("RGBA", (400, 240), (255, 255, 255, 255)).save(path)

            project = {
                "slide_width_inch": 13.333333,
                "image_width": 400,
                "image_height": 240,
                "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
                "pages": [
                    {
                        "page_no": 1,
                        "title": "第一页",
                        "summary": "摘要1",
                        "visual_image": str(visual_path_1),
                        "reference_image": str(reference_path_1),
                        "texts": [],
                    },
                    {
                        "page_no": 2,
                        "title": "第二页",
                        "summary": "摘要2",
                        "visual_image": str(visual_path_2),
                        "reference_image": str(reference_path_2),
                        "texts": [],
                    },
                ],
            }
            first_run_provider = FakeChatProvider(
                [
                    {"page_script": 'add_text(slide, "第一页成稿", 12, 14, 130, 36, size=20, color="163A63", bold=True)'},
                    {"page_script": 'add_text(slide, "第二页临时稿", 12, 14, 130, 36, size=20, color="163A63", bold=True)'},
                ]
            )
            second_run_provider = FakeChatProvider([])

            original_execute = execute_generated_text_script
            failing_state = {"triggered": False}

            def fail_after_first_page(
                script_path: Path,
                *,
                timeout_seconds: int = 600,
                stop_checker=None,
            ) -> None:
                if script_path.name == "generated_text_layout_preview_round_01.py":
                    original_execute(script_path, timeout_seconds=timeout_seconds, stop_checker=stop_checker)
                    if "page_02" in str(script_path):
                        failing_state["triggered"] = True
                        raise RuntimeError("模拟第 2 页导出失败")
                    return
                original_execute(script_path, timeout_seconds=timeout_seconds, stop_checker=stop_checker)

            with patch("ppt_system.export.direct_project_script.render_pptx_first_slide_to_png", return_value=None):
                with patch("ppt_system.export.direct_project_script.execute_generated_text_script", side_effect=fail_after_first_page):
                    with self.assertRaisesRegex(RuntimeError, "模拟第 2 页导出失败"):
                        export_project_to_pptx(
                            project,
                            work_dir,
                            output_pptx,
                            chat_provider=first_run_provider,  # type: ignore[arg-type]
                        )

            self.assertTrue(failing_state["triggered"])
            checkpoint_path = work_dir / "page_01" / CHECKPOINT_FILE_NAME
            self.assertTrue(checkpoint_path.exists())
            self.assertFalse((work_dir / "page_02" / CHECKPOINT_FILE_NAME).exists())
            self.assertTrue(list((work_dir / "page_02" / STEP_CHECKPOINT_DIR_NAME).glob("initial_script.*.json")))
            self.assertEqual(len(first_run_provider.calls), 2)

            with patch("ppt_system.export.direct_project_script.render_pptx_first_slide_to_png", return_value=None):
                result = export_project_to_pptx(
                    project,
                    work_dir,
                    output_pptx,
                    chat_provider=second_run_provider,  # type: ignore[arg-type]
                )

            self.assertTrue(output_pptx.exists())
            self.assertEqual(len(second_run_provider.calls), 0)
            self.assertIn("第二页临时稿", Path(result["text_script_path"]).read_text(encoding="utf-8"))
            self.assertIn("第一页成稿", Path(result["text_script_path"]).read_text(encoding="utf-8"))
            self.assertTrue((work_dir / "page_02" / CHECKPOINT_FILE_NAME).exists())


if __name__ == "__main__":
    unittest.main()
