from __future__ import annotations

import json
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

from ppt_system.composer import compose_pptx


class ComposerTests(unittest.TestCase):
    def test_compose_skips_empty_slot_text_and_invalid_asset_size(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            page_assets_dir = work_dir / "page_01" / "assets"
            page_assets_dir.mkdir(parents=True, exist_ok=True)

            Image.new("RGBA", (10, 10), (255, 0, 0, 255)).save(page_assets_dir / "asset_001.png")
            Image.new("RGBA", (10, 10), (0, 255, 0, 255)).save(page_assets_dir / "asset_002.png")
            (page_assets_dir / "assets.json").write_text(
                json.dumps(
                    {
                        "assets": [
                            {"file": "asset_001.png", "left": 10, "top": 10, "width": 300, "height": 200},
                            {"file": "asset_002.png", "left": 20, "top": 20, "width": 0, "height": 120},
                        ]
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )

            output_pptx = root / "result.pptx"
            compose_pptx(
                {
                    "slide_width_inch": 13.333333,
                    "image_width": 2048,
                    "image_height": 1152,
                    "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "000000"},
                    "pages": [
                        {
                            "page_no": 1,
                            "texts": [
                                {"text": "", "left": 0, "top": 0, "width": 500, "height": 120, "font_size": 0},
                                {"text": "有效标题", "left": 100, "top": 60, "width": 600, "height": 150, "font_size": 0},
                            ],
                        }
                    ],
                },
                work_dir,
                output_pptx,
            )

            self.assertTrue(output_pptx.exists())
            prs = Presentation(str(output_pptx))
            slide = prs.slides[0]
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            self.assertEqual(texts, ["有效标题"])
            text_shape = [shape for shape in slide.shapes if hasattr(shape, "text") and shape.text][0]
            self.assertEqual(text_shape.text_frame.auto_size, MSO_AUTO_SIZE.NONE)
            self.assertIn('wrap="none"', text_shape.element.xml)

    def test_compose_keeps_multiline_body_wrapped(self) -> None:
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            work_dir = root / "work"
            page_assets_dir = work_dir / "page_01" / "assets"
            page_assets_dir.mkdir(parents=True, exist_ok=True)
            (page_assets_dir / "assets.json").write_text(json.dumps({"assets": []}), encoding="utf-8")

            output_pptx = root / "result.pptx"
            compose_pptx(
                {
                    "slide_width_inch": 13.333333,
                    "image_width": 2048,
                    "image_height": 1152,
                    "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "000000"},
                    "pages": [
                        {
                            "page_no": 1,
                            "texts": [
                                {"text": "第一行\n第二行", "left": 100, "top": 60, "width": 600, "height": 180, "font_size": 24},
                            ],
                        }
                    ],
                },
                work_dir,
                output_pptx,
            )

            prs = Presentation(str(output_pptx))
            text_shape = [shape for shape in prs.slides[0].shapes if hasattr(shape, "text") and shape.text][0]
            self.assertTrue(text_shape.text_frame.word_wrap)


if __name__ == "__main__":
    unittest.main()
