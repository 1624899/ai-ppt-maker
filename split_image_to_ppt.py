from __future__ import annotations

import argparse
import json
from collections import deque
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation


EMU_PER_INCH = 914400


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="把透明 PNG 自动分割为多个独立 PNG，并按原位置还原到一页 PPT。"
    )
    parser.add_argument(
        "--input",
        default="ChatGPT Image5.png",
        help="输入 PNG 路径，默认使用当前目录下的 ChatGPT Image5.png。",
    )
    parser.add_argument(
        "--out-dir",
        default="split_elements",
        help="分割后 PNG 的输出目录。",
    )
    parser.add_argument(
        "--pptx",
        default="split_rebuild.pptx",
        help="输出 PPTX 文件路径。",
    )
    parser.add_argument(
        "--alpha-threshold",
        type=int,
        default=8,
        help="alpha 大于该值的像素会被视为元素像素。",
    )
    parser.add_argument(
        "--min-area",
        type=int,
        default=8,
        help="过滤过小噪点，单位为像素数量。",
    )
    parser.add_argument(
        "--padding",
        type=int,
        default=0,
        help="裁剪元素时额外保留的透明边距，单位为像素。",
    )
    parser.add_argument(
        "--slide-width-inch",
        type=float,
        default=13.333333,
        help="PPT 页面宽度，默认 16:9 宽屏。",
    )
    return parser.parse_args()


def find_components(mask: np.ndarray) -> list[dict[str, int | np.ndarray]]:
    height, width = mask.shape
    visited = np.zeros(mask.shape, dtype=bool)
    components: list[dict[str, int | np.ndarray]] = []

    # 用 8 邻域连通，避免抗锯齿线条在斜方向上被切得太碎。
    neighbors = (
        (-1, -1),
        (0, -1),
        (1, -1),
        (-1, 0),
        (1, 0),
        (-1, 1),
        (0, 1),
        (1, 1),
    )

    ys, xs = np.nonzero(mask)
    for start_x, start_y in zip(xs, ys):
        if visited[start_y, start_x]:
            continue

        queue: deque[tuple[int, int]] = deque([(start_x, start_y)])
        visited[start_y, start_x] = True
        pixels: list[tuple[int, int]] = []
        min_x = max_x = start_x
        min_y = max_y = start_y
        area = 0

        while queue:
            x, y = queue.popleft()
            pixels.append((x, y))
            area += 1
            if x < min_x:
                min_x = x
            elif x > max_x:
                max_x = x
            if y < min_y:
                min_y = y
            elif y > max_y:
                max_y = y

            for dx, dy in neighbors:
                nx = x + dx
                ny = y + dy
                if (
                    0 <= nx < width
                    and 0 <= ny < height
                    and mask[ny, nx]
                    and not visited[ny, nx]
                ):
                    visited[ny, nx] = True
                    queue.append((nx, ny))

        # 保存该连通域自己的像素坐标，后续裁剪时只保留这些像素，避免 bbox 内的小元素被重复裁入大元素。
        component_mask = np.zeros((max_y - min_y + 1, max_x - min_x + 1), dtype=bool)
        for x, y in pixels:
            component_mask[y - min_y, x - min_x] = True

        components.append(
            {
                "left": int(min_x),
                "top": int(min_y),
                "right": int(max_x + 1),
                "bottom": int(max_y + 1),
                "area": int(area),
                "mask": component_mask,
            }
        )

    return components


def save_components(
    image: Image.Image,
    components: list[dict[str, int | np.ndarray]],
    out_dir: Path,
    min_area: int,
    padding: int,
) -> list[dict[str, int | str]]:
    out_dir.mkdir(parents=True, exist_ok=True)
    width, height = image.size
    records: list[dict[str, int | str]] = []

    filtered = [component for component in components if int(component["area"]) >= min_area]
    filtered.sort(
        key=lambda item: (
            int(item["top"]),
            int(item["left"]),
            int(item["right"]) - int(item["left"]),
            int(item["bottom"]) - int(item["top"]),
        )
    )

    for index, component in enumerate(filtered, start=1):
        raw_left = int(component["left"])
        raw_top = int(component["top"])
        raw_right = int(component["right"])
        raw_bottom = int(component["bottom"])
        area = int(component["area"])
        component_mask = component["mask"]
        if not isinstance(component_mask, np.ndarray):
            raise TypeError("component mask must be a numpy array")

        left = max(0, raw_left - padding)
        top = max(0, raw_top - padding)
        right = min(width, raw_right + padding)
        bottom = min(height, raw_bottom + padding)

        crop = image.crop((left, top, right, bottom))
        crop_array = np.array(crop)

        # bbox 可能包含其他独立元素；这里只保留当前连通域自己的像素，其余全部设为透明。
        keep_mask = np.zeros((bottom - top, right - left), dtype=bool)
        mask_top = raw_top - top
        mask_left = raw_left - left
        keep_mask[
            mask_top : mask_top + component_mask.shape[0],
            mask_left : mask_left + component_mask.shape[1],
        ] = component_mask
        crop_array[~keep_mask, 3] = 0
        crop = Image.fromarray(crop_array, mode="RGBA")

        filename = f"element_{index:03d}.png"
        crop.save(out_dir / filename)

        records.append(
            {
                "index": index,
                "file": filename,
                "left": int(left),
                "top": int(top),
                "width": int(right - left),
                "height": int(bottom - top),
                "area": int(area),
            }
        )

    return records


def build_ppt(
    records: list[dict[str, int | str]],
    out_dir: Path,
    pptx_path: Path,
    image_width: int,
    image_height: int,
    slide_width_inch: float,
) -> None:
    prs = Presentation()
    prs.slide_width = int(slide_width_inch * EMU_PER_INCH)
    prs.slide_height = int(prs.slide_width * image_height / image_width)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    scale_x = prs.slide_width / image_width
    scale_y = prs.slide_height / image_height

    for record in records:
        file_path = out_dir / str(record["file"])
        left = int(int(record["left"]) * scale_x)
        top = int(int(record["top"]) * scale_y)
        width = int(int(record["width"]) * scale_x)
        height = int(int(record["height"]) * scale_y)
        slide.shapes.add_picture(str(file_path), left, top, width=width, height=height)

    prs.save(pptx_path)


def main() -> None:
    args = parse_args()
    input_path = Path(args.input)
    out_dir = Path(args.out_dir)
    pptx_path = Path(args.pptx)

    image = Image.open(input_path).convert("RGBA")
    alpha = np.array(image.getchannel("A"))
    mask = alpha > args.alpha_threshold

    components = find_components(mask)
    records = save_components(image, components, out_dir, args.min_area, args.padding)

    metadata = {
        "input": str(input_path),
        "image_width": image.width,
        "image_height": image.height,
        "alpha_threshold": args.alpha_threshold,
        "min_area": args.min_area,
        "padding": args.padding,
        "count": len(records),
        "elements": records,
    }
    (out_dir / "elements.json").write_text(
        json.dumps(metadata, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    build_ppt(
        records=records,
        out_dir=out_dir,
        pptx_path=pptx_path,
        image_width=image.width,
        image_height=image.height,
        slide_width_inch=args.slide_width_inch,
    )

    print(f"已保存 {len(records)} 个元素到: {out_dir}")
    print(f"已生成 PPT: {pptx_path}")


if __name__ == "__main__":
    main()
