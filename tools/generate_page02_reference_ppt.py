from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


BASE_DIR = Path(__file__).resolve().parents[1]
REFERENCE_IMAGE = BASE_DIR / "output" / "758b2965960a" / "01_reference_pages" / "page_02_reference.png"
ELEMENTS_IMAGE = BASE_DIR / "output" / "758b2965960a" / "02_elements_pages" / "page_02_elements.png"
OUTPUT_PPTX = BASE_DIR / "output" / "758b2965960a" / "page_02_reference_text_only.pptx"

SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)
IMAGE_WIDTH = 2048
IMAGE_HEIGHT = 1152

PRIMARY_BLUE = RGBColor(13, 37, 120)
ACCENT_BLUE = RGBColor(31, 93, 227)
WARNING_RED = RGBColor(255, 69, 51)


def px_to_emu(px: int, axis: str) -> Emu:
    total_px = IMAGE_WIDTH if axis == "x" else IMAGE_HEIGHT
    total_emu = SLIDE_WIDTH if axis == "x" else SLIDE_HEIGHT
    return Emu(round(px / total_px * total_emu))


def add_textbox(
    slide,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    font_size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    font_name: str = "Microsoft YaHei",
    align=PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_textbox(
        px_to_emu(left, "x"),
        px_to_emu(top, "y"),
        px_to_emu(width, "x"),
        px_to_emu(height, "y"),
    )
    shape.fill.background()
    shape.line.fill.background()

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return shape


def add_bullet_block(
    slide,
    items: list[str],
    left: int,
    top: int,
    width: int,
    line_gap: int,
    font_size: int,
    bullet_color: RGBColor,
    text_color: RGBColor,
):
    for index, item in enumerate(items):
        current_top = top + index * line_gap
        add_textbox(
            slide,
            "•",
            left,
            current_top,
            24,
            44,
            font_size + 4,
            bullet_color,
            bold=True,
            font_name="Arial",
        )
        add_textbox(
            slide,
            item,
            left + 38,
            current_top + 2,
            width - 38,
            44,
            font_size,
            text_color,
            font_name="Microsoft YaHei",
        )


def add_warning_label(slide):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        px_to_emu(100, "x"),
        px_to_emu(792, "y"),
        px_to_emu(206, "x"),
        px_to_emu(60, "y"),
    )
    box.fill.background()
    box.line.color.rgb = RGBColor(255, 92, 72)
    box.line.width = Pt(1)

    add_textbox(
        slide,
        "⚠",
        116,
        804,
        30,
        34,
        24,
        WARNING_RED,
        bold=True,
        font_name="Segoe UI Symbol",
    )
    add_textbox(
        slide,
        "效率风险",
        150,
        804,
        126,
        34,
        22,
        WARNING_RED,
        bold=True,
    )


def build_ppt():
    if not REFERENCE_IMAGE.exists():
        raise FileNotFoundError(f"未找到参考图: {REFERENCE_IMAGE}")
    if not ELEMENTS_IMAGE.exists():
        raise FileNotFoundError(f"未找到元素图: {ELEMENTS_IMAGE}")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture(str(ELEMENTS_IMAGE), 0, 0, width=SLIDE_WIDTH, height=SLIDE_HEIGHT)

    add_textbox(slide, "跳出聊天框", 73, 34, 610, 140, 50, PRIMARY_BLUE, bold=True, font_name="Microsoft YaHei")

    add_textbox(slide, "01", 112, 315, 80, 62, 32, RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "传统画图", 238, 319, 260, 56, 28, PRIMARY_BLUE, bold=True)

    add_bullet_block(
        slide,
        ["Draw.io 手动拉框", "对齐连线耗时", "流程维护成本高"],
        110,
        424,
        312,
        49,
        21,
        ACCENT_BLUE,
        PRIMARY_BLUE,
    )

    add_warning_label(slide)

    add_textbox(slide, "02", 792, 315, 80, 62, 32, RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "AI 转换", 917, 319, 220, 56, 28, PRIMARY_BLUE, bold=True)

    add_bullet_block(
        slide,
        ["输入业务逻辑文字", "指令：转化为 Draw.io XML", "AI 生成结构化代码"],
        789,
        424,
        392,
        49,
        21,
        ACCENT_BLUE,
        PRIMARY_BLUE,
    )

    add_textbox(slide, "AI", 1122, 693, 84, 70, 34, ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide, "03", 1431, 315, 80, 62, 32, RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "自动生成", 1554, 319, 240, 56, 28, PRIMARY_BLUE, bold=True)

    add_bullet_block(
        slide,
        ["复制 XML 到 Draw.io", "一秒生成整齐流程图", "结构化思维具象化"],
        1431,
        424,
        360,
        49,
        21,
        ACCENT_BLUE,
        PRIMARY_BLUE,
    )

    add_textbox(
        slide,
        "进阶玩法：让 AI 从聊天助手变成流程生产工具",
        462,
        976,
        1200,
        74,
        28,
        ACCENT_BLUE,
        bold=True,
    )

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))


if __name__ == "__main__":
    build_ppt()
