from __future__ import annotations

import json
from pathlib import Path
from tempfile import TemporaryDirectory
import sys

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from PIL import Image
from pptx import Presentation

from ppt_system.export_layer_mode import SEPARATE_LAYER_MODE
from ppt_system.text_script_runtime import build_project_script_source, execute_generated_text_script


def main() -> None:
    with TemporaryDirectory() as temp_dir:
        root = Path(temp_dir)
        work_dir = root / "work"
        page_assets_dir = work_dir / "page_01" / "assets"
        page_assets_dir.mkdir(parents=True, exist_ok=True)
        output_pptx = root / "layered_demo.pptx"

        Image.new("RGBA", (40, 30), (0, 128, 255, 255)).save(page_assets_dir / "asset_001.png")
        (page_assets_dir / "assets.json").write_text(
            json.dumps(
                {
                    "image_width": 400,
                    "image_height": 240,
                    "assets": [
                        {"index": 1, "file": "asset_001.png", "left": 40, "top": 50, "width": 120, "height": 80}
                    ],
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        project = {
            "slide_width_inch": 10.0,
            "image_width": 400,
            "image_height": 240,
            "default_font": {"font_name": "Microsoft YaHei", "font_size": 24, "color": "355C7D"},
            "pages": [
                {
                    "page_no": 1,
                    "title": "示例页",
                    "summary": "摘要",
                    "texts": [{"role": "title", "text": "示例页", "left": 40, "top": 50, "width": 120, "height": 80}],
                }
            ],
        }

        script_source = build_project_script_source(
            project,
            work_dir,
            output_pptx,
            [{"page_no": 1, "script": 'add_text_ref(slide, page_texts, "title", 40, 50, 120, 80, size=24, color="163A63", bold=True)'}],
            include_assets=True,
            layer_mode=SEPARATE_LAYER_MODE,
        )
        script_path = work_dir / "generated_text_layout.py"
        script_path.write_text(script_source, encoding="utf-8")
        execute_generated_text_script(script_path)

        prs = Presentation(str(output_pptx))
        summary = []
        for index, slide in enumerate(prs.slides, start=1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            picture_count = len([shape for shape in slide.shapes if shape.shape_type == 13])
            summary.append(
                {
                    "slide_no": index,
                    "picture_count": picture_count,
                    "texts": texts,
                }
            )
        print(
            json.dumps(
                {
                    "output_pptx": str(output_pptx),
                    "slide_count": len(prs.slides),
                    "slides": summary,
                },
                ensure_ascii=False,
                indent=2,
            )
        )


if __name__ == "__main__":
    main()
