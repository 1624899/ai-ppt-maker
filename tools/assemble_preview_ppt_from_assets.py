from __future__ import annotations

import argparse
import importlib.util
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="基于既有文字布局脚本与指定资产目录，组装预览 PPT。")
    parser.add_argument("--text-script", required=True, help="已有文字布局脚本路径。")
    parser.add_argument("--project", required=True, help="项目 JSON 路径，用于确定页面顺序。")
    parser.add_argument("--assets-root", required=True, help="按页资产根目录，例如 16_builtin_alpha_refine_edgeaware。")
    parser.add_argument("--output", required=True, help="输出 PPTX 路径。")
    parser.add_argument(
        "--pages",
        nargs="*",
        type=int,
        default=None,
        help="可选，显式指定页号；不传则按项目 JSON 中同时存在资产的页面导出。",
    )
    return parser.parse_args()


def load_module(script_path: Path):
    project_root = Path(__file__).resolve().parent.parent
    if str(project_root) not in sys.path:
        sys.path.insert(0, str(project_root))
    spec = importlib.util.spec_from_file_location("preview_text_layout_module", script_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"无法加载文字脚本：{script_path}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def resolve_page_numbers(project_path: Path, assets_root: Path, requested_pages: list[int] | None) -> list[int]:
    if requested_pages:
        return requested_pages

    project = json.loads(project_path.read_text(encoding="utf-8"))
    page_numbers: list[int] = []
    for page in project.get("pages", []):
        page_no = int(page.get("page_no", 0))
        if page_no <= 0:
            continue
        manifest_path = assets_root / f"page_{page_no:02d}" / "assets" / "assets.json"
        if manifest_path.exists():
            page_numbers.append(page_no)
    if not page_numbers:
        raise ValueError("没有找到可导出的页面：项目 JSON 与资产目录之间没有交集。")
    return page_numbers


def build_preview_ppt(
    text_script_path: Path,
    project_path: Path,
    assets_root: Path,
    output_path: Path,
    requested_pages: list[int] | None,
) -> Path:
    module = load_module(text_script_path)
    page_numbers = resolve_page_numbers(project_path, assets_root, requested_pages)

    prs = Presentation()
    prs.slide_width = Inches(float(module.SLIDE_W))
    prs.slide_height = Inches(float(module.SLIDE_H))
    blank_layout = prs.slide_layouts[6]

    for page_no in page_numbers:
        builder_name = f"build_slide_{page_no:02d}"
        builder = getattr(module, builder_name, None)
        if builder is None:
            raise AttributeError(f"文字脚本缺少页面构建函数：{builder_name}")

        manifest_path = assets_root / f"page_{page_no:02d}" / "assets" / "assets.json"
        if not manifest_path.exists():
            raise FileNotFoundError(f"第 {page_no} 页资产清单不存在：{manifest_path}")

        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        module.add_assets(slide, manifest_path)
        builder(slide)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> None:
    args = parse_args()
    output_path = build_preview_ppt(
        text_script_path=Path(args.text_script),
        project_path=Path(args.project),
        assets_root=Path(args.assets_root),
        output_path=Path(args.output),
        requested_pages=list(args.pages) if args.pages else None,
    )
    print(output_path.resolve())


if __name__ == "__main__":
    main()
