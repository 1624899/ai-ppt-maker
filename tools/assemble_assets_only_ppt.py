from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="基于已分割资产组装仅包含元素层的 PPT。")
    parser.add_argument("--project", required=True, help="项目 JSON 路径，用于读取页面尺寸与页序。")
    parser.add_argument("--assets-root", required=True, help="按页资产根目录，例如 output/<job_id>/03_ppt_build。")
    parser.add_argument("--output", required=True, help="输出 PPTX 路径。")
    parser.add_argument(
        "--pages",
        nargs="*",
        type=int,
        default=None,
        help="可选，显式指定页号；不传则自动导出存在 assets.json 的页面。",
    )
    return parser.parse_args()


def resolve_page_numbers(project_path: Path, assets_root: Path, requested_pages: list[int] | None) -> list[int]:
    if requested_pages:
        return [int(page_no) for page_no in requested_pages if int(page_no) > 0]

    project = json.loads(project_path.read_text(encoding="utf-8"))
    page_numbers: list[int] = []
    for page in project.get("pages", []):
        page_no = int(page.get("page_no", 0))
        if page_no <= 0:
            continue
        manifest_path = assets_root / f"page_{page_no:02d}" / "assets" / "assets.json"
        if manifest_path.exists():
            page_numbers.append(page_no)
    if not page_numbers:
        raise ValueError("没有找到可导出的页面：资产目录下不存在对应的 assets.json。")
    return page_numbers


def add_assets_to_slide(slide, manifest_path: Path, slide_width_inch: float, slide_height_inch: float) -> int:
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assets_dir = manifest_path.parent
    image_width = max(1, int(manifest.get("image_width", 1)))
    image_height = max(1, int(manifest.get("image_height", 1)))

    count = 0
    for asset in manifest.get("assets", []):
        asset_path = assets_dir / str(asset.get("file", "")).strip()
        if not asset_path.exists():
            continue
        width = int(asset.get("width", 0) or 0)
        height = int(asset.get("height", 0) or 0)
        if width <= 0 or height <= 0:
            continue
        left = int(asset.get("left", 0) or 0)
        top = int(asset.get("top", 0) or 0)
        slide.shapes.add_picture(
            str(asset_path),
            Inches(left / image_width * slide_width_inch),
            Inches(top / image_height * slide_height_inch),
            width=Inches(width / image_width * slide_width_inch),
            height=Inches(height / image_height * slide_height_inch),
        )
        count += 1
    return count


def build_assets_only_ppt(
    project_path: Path,
    assets_root: Path,
    output_path: Path,
    requested_pages: list[int] | None,
) -> dict[str, object]:
    project = json.loads(project_path.read_text(encoding="utf-8"))
    image_width = max(1, int(project.get("image_width", 2048)))
    image_height = max(1, int(project.get("image_height", 1152)))
    slide_width_inch = float(project.get("slide_width_inch", 13.333333))
    slide_height_inch = slide_width_inch * image_height / image_width
    page_numbers = resolve_page_numbers(project_path, assets_root, requested_pages)

    prs = Presentation()
    prs.slide_width = Inches(slide_width_inch)
    prs.slide_height = Inches(slide_height_inch)
    blank_layout = prs.slide_layouts[6]

    slides: list[dict[str, int]] = []
    for page_no in page_numbers:
        manifest_path = assets_root / f"page_{page_no:02d}" / "assets" / "assets.json"
        if not manifest_path.exists():
            raise FileNotFoundError(f"第 {page_no} 页资产清单不存在：{manifest_path}")
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        asset_count = add_assets_to_slide(slide, manifest_path, slide_width_inch, slide_height_inch)
        slides.append({"page_no": int(page_no), "asset_count": int(asset_count)})

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return {
        "output_pptx": str(output_path.resolve()),
        "slide_count": len(slides),
        "slides": slides,
    }


def main() -> None:
    args = parse_args()
    result = build_assets_only_ppt(
        project_path=Path(args.project),
        assets_root=Path(args.assets_root),
        output_path=Path(args.output),
        requested_pages=list(args.pages) if args.pages else None,
    )
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
