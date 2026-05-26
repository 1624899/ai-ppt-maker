from __future__ import annotations

import ast
import json
import runpy
from pathlib import Path
from typing import Any


CALL_CONTRACTS: dict[str, dict[str, Any]] = {
    "add_text": {"min_args": 6, "second_arg_kind": "literal"},
    "add_center_text": {"min_args": 6, "second_arg_kind": "literal"},
    "add_runs": {"min_args": 6, "second_arg_kind": "literal"},
    "add_text_ref": {"min_args": 7, "second_arg_kind": "page_texts"},
    "add_center_text_ref": {"min_args": 7, "second_arg_kind": "page_texts"},
}
ALLOWED_SCRIPT_CALLS = set(CALL_CONTRACTS)
ALLOWED_ALIGNS = {"LEFT", "CENTER", "RIGHT", "JUSTIFY"}
ALLOWED_ANCHORS = {"TOP", "MIDDLE", "BOTTOM"}
ASSET_RELATIVE_FIELDS = ("dx", "dy", "dw", "dh")
ASSET_ABSOLUTE_FIELDS = (
    ("left", "left"),
    ("top", "top"),
    ("width", "width"),
    ("height", "height"),
    ("x", "left"),
    ("y", "top"),
    ("w", "width"),
    ("h", "height"),
)


def execute_generated_text_script(script_path: Path, *, timeout_seconds: int = 600) -> None:
    namespace = runpy.run_path(str(script_path))
    build_deck = namespace.get("build_deck")
    if not callable(build_deck):
        raise RuntimeError(f"生成脚本缺少 build_deck 函数：{script_path}")
    output_path = Path(build_deck())
    if not output_path.exists():
        raise RuntimeError(f"执行生成的文字脚本后未发现输出文件：{output_path}")


def build_project_script_source(
    project: dict[str, Any],
    work_dir: Path,
    output_pptx: Path,
    page_scripts: list[dict[str, Any]],
    *,
    include_assets: bool = True,
) -> str:
    image_width = int(project.get("image_width", 2000))
    image_height = int(project.get("image_height", 1125))
    slide_width_inch = float(project.get("slide_width_inch", 13.333333))
    slide_height_inch = slide_width_inch * image_height / image_width
    default_font = project.get("default_font", {})
    font_name = str(default_font.get("font_name", "Microsoft YaHei"))
    font_color = str(default_font.get("color", "14254E")).lstrip("#").upper()
    page_texts = {
        int(page.get("page_no", 0)): _build_page_text_inventory(page)
        for page in project.get("pages", [])
        if int(page.get("page_no", 0)) > 0
    }
    page_asset_adjustments = {
        str(int(item["page_no"])): normalize_asset_adjustments(item.get("asset_adjustments"))
        for item in page_scripts
        if int(item.get("page_no", 0)) > 0 and normalize_asset_adjustments(item.get("asset_adjustments"))
    }

    page_functions = "\n\n".join(
        _build_page_function_source(item["page_no"], str(item["script"]))
        for item in page_scripts
    )
    page_builder_entries = ",\n        ".join(
        f"({int(item['page_no'])}, build_slide_{int(item['page_no']):02d})"
        for item in page_scripts
    )

    return f"""from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from ppt_system.text_style_runtime import should_wrap_text


# 基于参考图像素坐标映射到 16:9 PPT 页面坐标。
IMG_W = {image_width}
IMG_H = {image_height}
SLIDE_W = {slide_width_inch}
SLIDE_H = {slide_height_inch}
WORK_DIR = Path(r"{work_dir.resolve()}")
OUTPUT_PPTX = Path(r"{output_pptx.resolve()}")
DEFAULT_FONT_NAME = {font_name!r}
DEFAULT_FONT_COLOR = {font_color!r}
INCLUDE_ASSETS = {bool(include_assets)!r}
PAGE_TEXTS = {json.dumps(page_texts, ensure_ascii=False, indent=2)}
PAGE_ASSET_ADJUSTMENTS = {json.dumps(page_asset_adjustments, ensure_ascii=False, indent=2)}


def px_x(value):
    return Inches(value / IMG_W * SLIDE_W)


def px_y(value):
    return Inches(value / IMG_H * SLIDE_H)


def px_w(value):
    return Inches(value / IMG_W * SLIDE_W)


def px_h(value):
    return Inches(value / IMG_H * SLIDE_H)


def parse_color(value):
    color = str(value or DEFAULT_FONT_COLOR).strip().lstrip("#").upper()
    if len(color) != 6:
        color = DEFAULT_FONT_COLOR
    return RGBColor.from_string(color)


def parse_align(value):
    return getattr(PP_ALIGN, str(value or "LEFT").upper(), PP_ALIGN.LEFT)


def parse_anchor(value):
    return getattr(MSO_ANCHOR, str(value or "TOP").upper(), MSO_ANCHOR.TOP)


def _prepare_text_frame(shape, anchor="TOP", word_wrap=True):
    # 所有文字框均保持无填充、无线条，便于继续编辑或叠加到视觉稿。
    shape.fill.background()
    shape.line.fill.background()
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    text_frame.word_wrap = bool(word_wrap)
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = parse_anchor(anchor)
    return text_frame


def _apply_paragraph_style(paragraph, align="LEFT"):
    paragraph.alignment = parse_align(align)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05


def _apply_run_style(run, *, size=12, color=DEFAULT_FONT_COLOR, bold=False,
                     italic=False, font_name=DEFAULT_FONT_NAME):
    resolved_font_name = str(font_name or DEFAULT_FONT_NAME)
    run.font.name = resolved_font_name
    r_pr = run._r.get_or_add_rPr()
    r_pr.set(qn("a:ea"), resolved_font_name)
    r_pr.set(qn("a:cs"), resolved_font_name)
    latin = r_pr.get_or_add_latin()
    latin.typeface = resolved_font_name
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    run.font.color.rgb = parse_color(color)


def add_text(slide, text, x, y, w, h, size=12, color=DEFAULT_FONT_COLOR, bold=False,
             align="LEFT", font_name=DEFAULT_FONT_NAME, anchor="TOP", italic=False):
    shape = slide.shapes.add_textbox(px_x(x), px_y(y), px_w(w), px_h(h))
    text_frame = _prepare_text_frame(shape, anchor, should_wrap_text(text, w, h, size))

    lines = str(text).split("\\n")
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        _apply_paragraph_style(paragraph, align)
        run = paragraph.add_run()
        run.text = line
        _apply_run_style(
            run,
            size=size,
            color=color,
            bold=bold,
            italic=italic,
            font_name=font_name,
        )
    return shape


def add_center_text(slide, text, x, y, w, h, size=12, color=DEFAULT_FONT_COLOR, bold=False,
                    font_name=DEFAULT_FONT_NAME, anchor="MIDDLE", italic=False):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align="CENTER",
        font_name=font_name,
        anchor=anchor,
        italic=italic,
    )


def add_runs(slide, runs, x, y, w, h, align="LEFT", font_name=DEFAULT_FONT_NAME, anchor="TOP"):
    shape = slide.shapes.add_textbox(px_x(x), px_y(y), px_w(w), px_h(h))
    text_frame = _prepare_text_frame(shape, anchor, False)
    paragraph = text_frame.paragraphs[0]
    _apply_paragraph_style(paragraph, align)
    for item in list(runs or []):
        run = paragraph.add_run()
        run.text = str(item.get("text", ""))
        _apply_run_style(
            run,
            size=float(item.get("size", 12)),
            color=item.get("color", DEFAULT_FONT_COLOR),
            bold=bool(item.get("bold", False)),
            italic=bool(item.get("italic", False)),
            font_name=item.get("font_name", font_name),
        )
    return shape


def add_text_ref(slide, page_texts, text_id, x, y, w, h, size=12, color=DEFAULT_FONT_COLOR, bold=False,
                 align="LEFT", font_name=DEFAULT_FONT_NAME, anchor="TOP", italic=False):
    text = str(page_texts.get(str(text_id), ""))
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=align,
        font_name=font_name,
        anchor=anchor,
        italic=italic,
    )


def add_center_text_ref(slide, page_texts, text_id, x, y, w, h, size=12, color=DEFAULT_FONT_COLOR, bold=False,
                        font_name=DEFAULT_FONT_NAME, anchor="MIDDLE", italic=False):
    text = str(page_texts.get(str(text_id), ""))
    return add_center_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        font_name=font_name,
        anchor=anchor,
        italic=italic,
    )


def _resolve_asset_box(asset, adjustment_plan):
    left = int(asset["left"])
    top = int(asset["top"])
    width = int(asset["width"])
    height = int(asset["height"])

    global_adjustment = dict((adjustment_plan or {{}}).get("global", {{}}))
    left += int(global_adjustment.get("dx", 0))
    top += int(global_adjustment.get("dy", 0))
    width += int(global_adjustment.get("dw", 0))
    height += int(global_adjustment.get("dh", 0))

    asset_map = dict((adjustment_plan or {{}}).get("asset_map", {{}}))
    asset_adjustment = dict(asset_map.get(str(int(asset.get("index", 0))), {{}}))
    if "left" in asset_adjustment:
        left = int(asset_adjustment["left"])
    else:
        left += int(asset_adjustment.get("dx", 0))
    if "top" in asset_adjustment:
        top = int(asset_adjustment["top"])
    else:
        top += int(asset_adjustment.get("dy", 0))
    if "width" in asset_adjustment:
        width = int(asset_adjustment["width"])
    else:
        width += int(asset_adjustment.get("dw", 0))
    if "height" in asset_adjustment:
        height = int(asset_adjustment["height"])
    else:
        height += int(asset_adjustment.get("dh", 0))
    return left, top, width, height


def add_assets(slide, manifest_path, page_no):
    manifest = json.loads(Path(manifest_path).read_text(encoding="utf-8"))
    assets_dir = Path(manifest_path).parent
    asset_img_w = max(1, int(manifest.get("image_width", IMG_W) or IMG_W))
    asset_img_h = max(1, int(manifest.get("image_height", IMG_H) or IMG_H))
    adjustment_plan = PAGE_ASSET_ADJUSTMENTS.get(str(page_no), {{}})
    for asset in manifest.get("assets", []):
        asset_path = assets_dir / str(asset["file"])
        left, top, width, height = _resolve_asset_box(asset, adjustment_plan)
        if width <= 0 or height <= 0:
            continue
        slide.shapes.add_picture(
            str(asset_path),
            Inches(left / asset_img_w * SLIDE_W),
            Inches(top / asset_img_h * SLIDE_H),
            width=Inches(width / asset_img_w * SLIDE_W),
            height=Inches(height / asset_img_h * SLIDE_H),
        )


{page_functions}


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]
    page_builders = [
        {page_builder_entries}
    ]
    for page_no, builder in page_builders:
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        if INCLUDE_ASSETS:
            add_assets(slide, WORK_DIR / f"page_{{page_no:02d}}" / "assets" / "assets.json", page_no)
        builder(slide)
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


if __name__ == "__main__":
    path = build_deck()
    print(path.resolve())
"""


def normalize_page_script(script: str) -> str:
    normalized_lines: list[str] = []
    for raw_line in _coalesce_script_lines(str(script)):
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            normalized_lines.append("")
            continue
        if stripped.startswith("#"):
            normalized_lines.append(stripped)
            continue
        sanitized = _sanitize_script_line(stripped)
        _validate_script_call(sanitized)
        stripped = sanitized
        normalized_lines.append(stripped)
    while normalized_lines and not normalized_lines[-1]:
        normalized_lines.pop()
    return "\n".join(normalized_lines)


def normalize_asset_adjustments(adjustments: Any) -> dict[str, Any]:
    if not isinstance(adjustments, dict):
        return {}

    normalized: dict[str, Any] = {}
    global_adjustment = _normalize_numeric_mapping(adjustments.get("global"), ASSET_RELATIVE_FIELDS)
    if global_adjustment:
        normalized["global"] = global_adjustment

    asset_map: dict[str, dict[str, int]] = {}
    raw_asset_map = adjustments.get("asset_map")
    if isinstance(raw_asset_map, dict):
        for raw_index, raw_plan in raw_asset_map.items():
            index = _coerce_positive_int(raw_index)
            if index is None:
                continue
            plan = _normalize_single_asset_adjustment(raw_plan)
            if plan:
                asset_map[str(index)] = plan

    raw_assets = adjustments.get("assets")
    if isinstance(raw_assets, list):
        for item in raw_assets:
            if not isinstance(item, dict):
                continue
            index = _coerce_positive_int(item.get("index"))
            if index is None:
                continue
            plan = _normalize_single_asset_adjustment(item)
            if plan:
                asset_map[str(index)] = plan

    if asset_map:
        normalized["asset_map"] = asset_map
    return normalized


def _validate_script_call(line: str) -> None:
    node = ast.parse(line, mode="exec")
    if len(node.body) != 1 or not isinstance(node.body[0], ast.Expr):
        raise RuntimeError(f"脚本行不合法：{line}")
    expression = node.body[0].value
    if not isinstance(expression, ast.Call) or not isinstance(expression.func, ast.Name):
        raise RuntimeError(f"脚本调用不合法：{line}")

    function_name = expression.func.id
    contract = CALL_CONTRACTS.get(function_name)
    if contract is None:
        raise RuntimeError(f"脚本调用超出白名单：{function_name}")
    if len(expression.args) < int(contract["min_args"]):
        raise RuntimeError(f"脚本参数不足：{line}")

    first_arg = expression.args[0]
    if not isinstance(first_arg, ast.Name) or first_arg.id != "slide":
        raise RuntimeError(f"第一个参数必须是 slide：{line}")

    second_arg_kind = str(contract["second_arg_kind"])
    if second_arg_kind == "page_texts":
        second_arg = expression.args[1]
        if not isinstance(second_arg, ast.Name) or second_arg.id != "page_texts":
            raise RuntimeError(f"第二个参数必须是 page_texts：{line}")
        literal_args = expression.args[2:]
    else:
        literal_args = expression.args[1:]

    for arg in literal_args:
        _literal_eval(arg)
    for keyword in expression.keywords:
        if keyword.arg is None:
            raise RuntimeError(f"不允许使用 **kwargs：{line}")
        value = _literal_eval(keyword.value)
        if keyword.arg == "align":
            align = str(value).upper()
            if align not in ALLOWED_ALIGNS:
                raise RuntimeError(f"align 不合法：{align}")
        if keyword.arg == "anchor":
            anchor = str(value).upper()
            if anchor not in ALLOWED_ANCHORS:
                raise RuntimeError(f"anchor 不合法：{anchor}")


def _sanitize_script_line(line: str) -> str:
    """对模型偶发返回的损坏转义做最小修复，避免整次导出中断。"""
    sanitized = str(line).replace("\\\r", "\\r").replace("\\\n", "\\n")
    if sanitized.count('"') % 2 != 0:
        sanitized = sanitized.replace("\\", "\\\\")
    return sanitized


def _coalesce_script_lines(script: str) -> list[str]:
    """把模型返回的多行函数调用拼回单行，再进入 AST 校验。"""
    result: list[str] = []
    buffer: list[str] = []
    paren_depth = 0

    for raw_line in str(script).splitlines():
        stripped = raw_line.strip()
        if not stripped:
            if not buffer:
                result.append("")
            continue
        if stripped.startswith("#") and not buffer:
            result.append(stripped)
            continue

        buffer.append(stripped)
        paren_depth += stripped.count("(") - stripped.count(")")
        if paren_depth > 0:
            continue

        result.append(" ".join(buffer))
        buffer = []
        paren_depth = 0

    if buffer:
        result.append(" ".join(buffer))
    return result


def _literal_eval(node: ast.AST) -> Any:
    try:
        return ast.literal_eval(node)
    except Exception as exc:
        raise RuntimeError(f"脚本参数必须是字面量：{ast.dump(node)}") from exc


def _build_page_function_source(page_no: int, script: str) -> str:
    body = script or "pass"
    page_texts_line = f'page_texts = PAGE_TEXTS["{page_no}"]'
    full_body = "\n".join([page_texts_line, body]) if body != "pass" else f"{page_texts_line}\npass"
    indented = "\n".join(f"    {line}" if line else "" for line in full_body.splitlines()) or "    pass"
    return f"def build_slide_{page_no:02d}(slide):\n{indented}"


def _build_page_text_inventory(page: dict[str, Any]) -> dict[str, str]:
    inventory: dict[str, str] = {}
    title = str(page.get("title", "")).strip()
    if title:
        inventory["title"] = title

    bullets = page.get("bullets", [])
    if isinstance(bullets, list):
        normalized_bullets = [str(item).strip() for item in bullets if str(item).strip()]
        for index, bullet in enumerate(normalized_bullets, start=1):
            inventory[f"bullet_{index}"] = bullet

    if not inventory.get("bullet_1"):
        body_index = 1
        for item in page.get("texts", []):
            if not isinstance(item, dict):
                continue
            text = str(item.get("text", "")).strip()
            role = str(item.get("role", "")).strip().lower()
            if not text or role == "title":
                continue
            inventory[f"body_{body_index}"] = text
            body_index += 1

    summary = str(page.get("summary", "")).strip()
    if summary:
        inventory["summary"] = summary

    for item in page.get("texts", []):
        if not isinstance(item, dict):
            continue
        text = str(item.get("text", "")).strip()
        if not text:
            continue
        role = str(item.get("role", "")).strip().lower() or "text"
        key_prefix = role if role not in inventory else f"{role}_text"
        next_index = 1
        candidate_key = f"{key_prefix}_{next_index}"
        while candidate_key in inventory:
            next_index += 1
            candidate_key = f"{key_prefix}_{next_index}"
        inventory[candidate_key] = text

    if not inventory:
        inventory["summary"] = summary
    return inventory


def _normalize_single_asset_adjustment(raw_plan: Any) -> dict[str, int]:
    if not isinstance(raw_plan, dict):
        return {}

    normalized = _normalize_numeric_mapping(raw_plan, ASSET_RELATIVE_FIELDS)
    for source_name, target_name in ASSET_ABSOLUTE_FIELDS:
        if target_name in normalized:
            continue
        value = _coerce_int(raw_plan.get(source_name))
        if value is None:
            continue
        if target_name in {"width", "height"} and value <= 0:
            continue
        normalized[target_name] = value
    return normalized


def _normalize_numeric_mapping(raw_mapping: Any, fields: tuple[str, ...]) -> dict[str, int]:
    if not isinstance(raw_mapping, dict):
        return {}

    normalized: dict[str, int] = {}
    for field in fields:
        value = _coerce_int(raw_mapping.get(field))
        if value is None or value == 0:
            continue
        normalized[field] = value
    return normalized


def _coerce_positive_int(value: Any) -> int | None:
    result = _coerce_int(value)
    if result is None or result <= 0:
        return None
    return result


def _coerce_int(value: Any) -> int | None:
    try:
        return int(round(float(value)))
    except (TypeError, ValueError):
        return None
