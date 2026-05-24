from __future__ import annotations

import json
import unicodedata
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from ppt_system.composer_runtime import (
    resolve_asset_geometry,
    resolve_font_size_pt,
    resolve_text_geometry,
    should_render_text_item,
)
from ppt_system.text_style_runtime import apply_run_font_family, scale_font_size_pt


EMU_PER_INCH = 914400


def _emu(value: float) -> int:
    return int(value)


def add_text_box(slide, text_item: dict[str, Any], scale_x: float, scale_y: float, default_font: dict[str, Any]) -> None:
    geometry = resolve_text_geometry(text_item, scale_x, scale_y)
    if geometry is None or not should_render_text_item(text_item):
        return

    box = slide.shapes.add_textbox(
        geometry.left,
        geometry.top,
        geometry.width,
        geometry.height,
    )

    # 文字框无填充、无线条，保证只留下可编辑文字。
    box.fill.background()
    box.line.fill.background()

    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = _should_wrap_text(text_item, default_font)
    # 保持实际导出的 PPT 与脚本设定字号一致，不交给 Office 自动缩放。
    frame.auto_size = MSO_AUTO_SIZE.NONE

    paragraph = frame.paragraphs[0]
    paragraph.alignment = getattr(PP_ALIGN, str(text_item.get("align", default_font.get("align", "LEFT"))).upper())
    run = paragraph.add_run()
    run.text = str(text_item.get("text", ""))

    font = run.font
    resolved_font_name = str(text_item.get("font_name", default_font.get("font_name", "Microsoft YaHei")))
    apply_run_font_family(run, resolved_font_name)
    font_scale = float(default_font.get("render_font_scale", 1.0) or 1.0)
    font.size = Pt(scale_font_size_pt(resolve_font_size_pt(text_item, default_font), scale=font_scale))
    font.bold = bool(text_item.get("bold", default_font.get("bold", False)))
    font.italic = bool(text_item.get("italic", default_font.get("italic", False)))
    color = str(text_item.get("color", default_font.get("color", "FFFFFF"))).lstrip("#")
    font.color.rgb = RGBColor.from_string(color)


def _should_wrap_text(text_item: dict[str, Any], default_font: dict[str, Any]) -> bool:
    text = str(text_item.get("text", ""))
    if "\n" in text:
        return True
    try:
        font_size = resolve_font_size_pt(text_item, default_font)
        width = float(text_item.get("width", 0))
        if _estimate_text_width(text, font_size) <= width * 1.05:
            return False
        return float(text_item.get("height", 0)) > font_size * 2.4
    except (TypeError, ValueError):
        return True


def _estimate_text_width(text: str, font_size: float) -> float:
    units = 0.0
    for char in str(text):
        if char.isspace():
            units += 0.35
        elif unicodedata.east_asian_width(char) in {"F", "W"}:
            units += 1.0
        else:
            units += 0.55
    return units * float(font_size)


def compose_pptx(project: dict[str, Any], work_dir: Path, output_pptx: Path) -> None:
    prs = Presentation()
    slide_width_inch = float(project.get("slide_width_inch", 13.333333))
    image_width = int(project.get("image_width", 2000))
    image_height = int(project.get("image_height", 1125))
    prs.slide_width = int(slide_width_inch * EMU_PER_INCH)
    prs.slide_height = int(prs.slide_width * image_height / image_width)

    scale_x = prs.slide_width / image_width
    scale_y = prs.slide_height / image_height
    default_font = project.get("default_font", {})

    for page in project["pages"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        page_dir = work_dir / f"page_{int(page['page_no']):02d}"
        assets = json.loads((page_dir / "assets" / "assets.json").read_text(encoding="utf-8"))

        for asset in assets["assets"]:
            asset_path = page_dir / "assets" / str(asset["file"])
            geometry = resolve_asset_geometry(asset, scale_x, scale_y)
            if geometry is None:
                continue
            slide.shapes.add_picture(
                str(asset_path),
                geometry.left,
                geometry.top,
                width=geometry.width,
                height=geometry.height,
            )

        for text_item in page.get("texts", []):
            add_text_box(slide, text_item, scale_x, scale_y, default_font)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
