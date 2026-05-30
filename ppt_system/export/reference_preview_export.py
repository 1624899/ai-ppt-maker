from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches


def export_reference_images_to_pptx(
    reference_pages: list[dict[str, Any]],
    job_dir: Path,
    output_pptx: Path,
    *,
    image_width: int,
    image_height: int,
) -> dict[str, Any]:
    if not reference_pages:
        raise ValueError("缺少参考图，无法导出图片版 PPT。")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(13.333333 * image_height / max(1, image_width))
    blank_layout = prs.slide_layouts[6]

    for item in sorted(reference_pages, key=lambda page: int(page.get("page_no", 0))):
        image_ref = str(item.get("image", "")).strip()
        if not image_ref:
            raise ValueError(f"第 {item.get('page_no', '?')} 页参考图路径为空。")
        image_path = resolve_job_image_path(job_dir, image_ref)
        if not image_path.exists():
            raise FileNotFoundError(f"参考图不存在：{image_path}")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return {
        "pptx_path": str(output_pptx),
        "page_count": len(reference_pages),
    }


def resolve_job_image_path(job_dir: Path, image_ref: str) -> Path:
    normalized = str(image_ref or "").strip()
    if not normalized:
        raise ValueError("图像路径为空。")
    candidate = Path(normalized)
    if candidate.is_absolute():
        return candidate
    stripped = normalized.lstrip("/\\")
    parts = Path(stripped).parts
    if len(parts) >= 3 and parts[0] == "runs":
        return job_dir / Path(*parts[2:])
    return job_dir / stripped
