from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path
from typing import Any

from ppt_system.export.export_layer_mode import OVERLAY_LAYER_MODE, build_slide_layer_specs, normalize_layer_mode
from ppt_system.export.text_script_schema import normalize_page_script


ASSET_RELATIVE_FIELDS = ("dx", "dy", "dw", "dh")
ASSET_ABSOLUTE_FIELDS = (
    ("left", "left"),
    ("top", "top"),
    ("width", "width"),
    ("height", "height"),
    ("x", "left"),
    ("y", "top"),
    ("w", "width"),
    ("h", "height"),
)


def execute_generated_text_script(script_path: Path, *, timeout_seconds: int = 600) -> None:
    worker_script = Path(__file__).with_name("text_script_worker.py")
    command = [sys.executable, str(worker_script), str(Path(script_path).resolve())]
    try:
        completed = subprocess.run(
            command,
            capture_output=True,
            text=False,
            timeout=max(1, int(timeout_seconds)),
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        stdout = _decode_subprocess_output(exc.stdout).strip()
        stderr = _decode_subprocess_output(exc.stderr).strip()
        detail_parts = [f"生成脚本执行超时：{script_path}，超过 {int(timeout_seconds)} 秒"]
        if stdout:
            detail_parts.append(f"stdout:\n{stdout}")
        if stderr:
            detail_parts.append(f"stderr:\n{stderr}")
        raise TimeoutError("\n".join(detail_parts)) from exc

    stdout = _decode_subprocess_output(completed.stdout).strip()
    stderr = _decode_subprocess_output(completed.stderr).strip()
    payload = _load_worker_payload(stdout, stderr, script_path)
    if completed.returncode != 0:
        detail_parts = [f"生成脚本执行失败：{script_path}，退出码 {completed.returncode}"]
        if isinstance(payload, dict):
            error = str(payload.get("error", "")).strip()
            traceback_text = str(payload.get("traceback", "")).strip()
            script_stdout = str(payload.get("script_stdout", "")).strip()
            script_stderr = str(payload.get("script_stderr", "")).strip()
            if error:
                detail_parts.append(f"error:\n{error}")
            if traceback_text:
                detail_parts.append(f"traceback:\n{traceback_text}")
            if script_stdout:
                detail_parts.append(f"script_stdout:\n{script_stdout}")
            if script_stderr:
                detail_parts.append(f"script_stderr:\n{script_stderr}")
        elif stdout:
            detail_parts.append(f"stdout:\n{stdout}")
        if stderr:
            detail_parts.append(f"stderr:\n{stderr}")
        raise RuntimeError("\n".join(detail_parts))

    if not bool(payload.get("ok")):
        raise RuntimeError(f"生成脚本执行返回失败结果：{payload}")
    output_path = Path(str(payload.get("output_path", "")).strip())
    if not output_path.exists():
        raise RuntimeError(f"执行生成的文字脚本后未发现输出文件：{output_path}")


def _decode_subprocess_output(raw_output: Any) -> str:
    if raw_output is None:
        return ""
    if isinstance(raw_output, str):
        return raw_output
    if isinstance(raw_output, bytes):
        return raw_output.decode("utf-8", errors="replace")
    return str(raw_output)


def _load_worker_payload(stdout: str, stderr: str, script_path: Path) -> dict[str, Any]:
    if not stdout:
        detail_parts = [f"生成脚本执行后没有返回结果：{script_path}"]
        if stderr:
            detail_parts.append(f"stderr:\n{stderr}")
        raise RuntimeError("\n".join(detail_parts))
    try:
        payload = json.loads(stdout)
    except json.JSONDecodeError as exc:
        detail_parts = [f"生成脚本执行返回了非 JSON 输出：{script_path}", f"stdout:\n{stdout}"]
        if stderr:
            detail_parts.append(f"stderr:\n{stderr}")
        raise RuntimeError("\n".join(detail_parts)) from exc
    if not isinstance(payload, dict):
        raise RuntimeError(f"生成脚本执行返回结果不是对象：{payload!r}")
    return payload


def build_project_script_source(
    project: dict[str, Any],
    work_dir: Path,
    output_pptx: Path,
    page_scripts: list[dict[str, Any]],
    *,
    include_assets: bool = True,
    layer_mode: str = OVERLAY_LAYER_MODE,
) -> str:
    image_width = int(project.get("image_width", 2000))
    image_height = int(project.get("image_height", 1125))
    slide_width_inch = float(project.get("slide_width_inch", 13.333333))
    slide_height_inch = slide_width_inch * image_height / image_width
    default_font = project.get("default_font", {})
    font_name = str(default_font.get("font_name", "Microsoft YaHei"))
    font_color = str(default_font.get("color", "14254E")).lstrip("#").upper()
    page_texts = {
        int(page.get("page_no", 0)): _build_page_text_inventory(page)
        for page in project.get("pages", [])
        if int(page.get("page_no", 0)) > 0
    }
    page_asset_adjustments = {
        str(int(item["page_no"])): normalize_asset_adjustments(item.get("asset_adjustments"))
        for item in page_scripts
        if int(item.get("page_no", 0)) > 0 and normalize_asset_adjustments(item.get("asset_adjustments"))
    }
    slide_layer_specs = [spec.to_payload() for spec in build_slide_layer_specs(layer_mode)]
    resolved_layer_mode = normalize_layer_mode(layer_mode)

    page_functions = "\n\n".join(
        _build_page_function_source(item["page_no"], str(item["script"]))
        for item in page_scripts
    )
    page_builder_entries = ",\n        ".join(
        f"({int(item['page_no'])}, build_slide_{int(item['page_no']):02d})"
        for item in page_scripts
    )

    return f"""from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from ppt_system.export.text_style_runtime import should_wrap_text


# 基于原稿图像素坐标映射到 16:9 PPT 页面坐标。
IMG_W = {image_width}
IMG_H = {image_height}
SLIDE_W = {slide_width_inch}
SLIDE_H = {slide_height_inch}
WORK_DIR = Path(r"{work_dir.resolve()}")
OUTPUT_PPTX = Path(r"{output_pptx.resolve()}")
DEFAULT_FONT_NAME = {font_name!r}
DEFAULT_FONT_COLOR = {font_color!r}
INCLUDE_ASSETS = {bool(include_assets)!r}
LAYER_MODE = {resolved_layer_mode!r}
PAGE_TEXTS = {json.dumps(page_texts, ensure_ascii=False, indent=2)}
PAGE_ASSET_ADJUSTMENTS = {json.dumps(page_asset_adjustments, ensure_ascii=False, indent=2)}
SLIDE_LAYER_SPECS = {slide_layer_specs!r}


def px_x(value):
    return Inches(value / IMG_W * SLIDE_W)


def px_y(value):
    return Inches(value / IMG_H * SLIDE_H)


def px_w(value):
    return Inches(value / IMG_W * SLIDE_W)


def px_h(value):
    return Inches(value / IMG_H * SLIDE_H)


def parse_color(value):
    color = str(value or DEFAULT_FONT_COLOR).strip().lstrip("#").upper()
    if len(color) != 6:
        color = DEFAULT_FONT_COLOR
    return RGBColor.from_string(color)


def parse_align(value):
    return getattr(PP_ALIGN, str(value or "LEFT").upper(), PP_ALIGN.LEFT)


def parse_anchor(value):
    return getattr(MSO_ANCHOR, str(value or "TOP").upper(), MSO_ANCHOR.TOP)


def _prepare_text_frame(shape, anchor="TOP", word_wrap=True):
    # 所有文字框均保持无填充、无线条，便于继续编辑或叠加到视觉稿。
    shape.fill.background()
    shape.line.fill.background()
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    text_frame.word_wrap = bool(word_wrap)
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = parse_anchor(anchor)
    return text_frame


def _apply_paragraph_style(paragraph, align="LEFT"):
    paragraph.alignment = parse_align(align)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05


def _apply_run_style(run, *, size=12, color=DEFAULT_FONT_COLOR, bold=False,
                     italic=False, font_name=DEFAULT_FONT_NAME):
    resolved_font_name = str(font_name or DEFAULT_FONT_NAME)
    run.font.name = resolved_font_name
    r_pr = run._r.get_or_add_rPr()
    r_pr.set(qn("a:ea"), resolved_font_name)
    r_pr.set(qn("a:cs"), resolved_font_name)
    latin = r_pr.get_or_add_latin()
    latin.typeface = resolved_font_name
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.italic = bool(italic)
    run.font.color.rgb = parse_color(color)


def add_text(slide, text, x, y, w, h, size=12, color=DEFAULT_FONT_COLOR, bold=False,
             align="LEFT", font_name=DEFAULT_FONT_NAME, anchor="TOP", italic=False):
    shape = slide.shapes.add_textbox(px_x(x), px_y(y), px_w(w), px_h(h))
    text_frame = _prepare_text_frame(shape, anchor, should_wrap_text(text, w, h, size))

    lines = str(text).split("\\n")
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        _apply_paragraph_style(paragraph, align)
        run = paragraph.add_run()
        run.text = line
        _apply_run_style(
            run,
            size=size,
            color=color,
            bold=bold,
            italic=italic,
            font_name=font_name,
        )
    return shape


def add_center_text(slide, text, x, y, w, h, size=12, color=DEFAULT_FONT_COLOR, bold=False,
                    font_name=DEFAULT_FONT_NAME, anchor="MIDDLE", italic=False):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align="CENTER",
        font_name=font_name,
        anchor=anchor,
        italic=italic,
    )


def add_runs(slide, runs, x, y, w, h, align="LEFT", font_name=DEFAULT_FONT_NAME, anchor="TOP"):
    shape = slide.shapes.add_textbox(px_x(x), px_y(y), px_w(w), px_h(h))
    text_frame = _prepare_text_frame(shape, anchor, False)
    paragraph = text_frame.paragraphs[0]
    _apply_paragraph_style(paragraph, align)
    for item in list(runs or []):
        run = paragraph.add_run()
        run.text = str(item.get("text", ""))
        _apply_run_style(
            run,
            size=float(item.get("size", 12)),
            color=item.get("color", DEFAULT_FONT_COLOR),
            bold=bool(item.get("bold", False)),
            italic=bool(item.get("italic", False)),
            font_name=item.get("font_name", font_name),
        )
    return shape


def add_text_ref(slide, page_texts, text_id, x, y, w, h, size=12, color=DEFAULT_FONT_COLOR, bold=False,
                 align="LEFT", font_name=DEFAULT_FONT_NAME, anchor="TOP", italic=False):
    text = str(page_texts.get(str(text_id), ""))
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        align=align,
        font_name=font_name,
        anchor=anchor,
        italic=italic,
    )


def add_center_text_ref(slide, page_texts, text_id, x, y, w, h, size=12, color=DEFAULT_FONT_COLOR, bold=False,
                        font_name=DEFAULT_FONT_NAME, anchor="MIDDLE", italic=False):
    text = str(page_texts.get(str(text_id), ""))
    return add_center_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        bold=bold,
        font_name=font_name,
        anchor=anchor,
        italic=italic,
    )


def _resolve_asset_box(asset, adjustment_plan):
    left = int(asset["left"])
    top = int(asset["top"])
    width = int(asset["width"])
    height = int(asset["height"])

    global_adjustment = dict((adjustment_plan or {{}}).get("global", {{}}))
    left += int(global_adjustment.get("dx", 0))
    top += int(global_adjustment.get("dy", 0))
    width += int(global_adjustment.get("dw", 0))
    height += int(global_adjustment.get("dh", 0))

    asset_map = dict((adjustment_plan or {{}}).get("asset_map", {{}}))
    asset_adjustment = dict(asset_map.get(str(int(asset.get("index", 0))), {{}}))
    if "left" in asset_adjustment:
        left = int(asset_adjustment["left"])
    else:
        left += int(asset_adjustment.get("dx", 0))
    if "top" in asset_adjustment:
        top = int(asset_adjustment["top"])
    else:
        top += int(asset_adjustment.get("dy", 0))
    if "width" in asset_adjustment:
        width = int(asset_adjustment["width"])
    else:
        width += int(asset_adjustment.get("dw", 0))
    if "height" in asset_adjustment:
        height = int(asset_adjustment["height"])
    else:
        height += int(asset_adjustment.get("dh", 0))
    return left, top, width, height


def add_assets(slide, manifest_path, page_no):
    manifest = json.loads(Path(manifest_path).read_text(encoding="utf-8"))
    assets_dir = Path(manifest_path).parent
    asset_img_w = max(1, int(manifest.get("image_width", IMG_W) or IMG_W))
    asset_img_h = max(1, int(manifest.get("image_height", IMG_H) or IMG_H))
    adjustment_plan = PAGE_ASSET_ADJUSTMENTS.get(str(page_no), {{}})
    for asset in manifest.get("assets", []):
        asset_path = assets_dir / str(asset["file"])
        left, top, width, height = _resolve_asset_box(asset, adjustment_plan)
        if width <= 0 or height <= 0:
            continue
        slide.shapes.add_picture(
            str(asset_path),
            Inches(left / asset_img_w * SLIDE_W),
            Inches(top / asset_img_h * SLIDE_H),
            width=Inches(width / asset_img_w * SLIDE_W),
            height=Inches(height / asset_img_h * SLIDE_H),
        )


{page_functions}


def add_page_content(slide, page_no, builder, layer_spec):
    if bool(layer_spec.get("include_assets")) and INCLUDE_ASSETS:
        add_assets(slide, WORK_DIR / f"page_{{page_no:02d}}" / "assets" / "assets.json", page_no)
    if bool(layer_spec.get("include_text")):
        builder(slide)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]
    page_builders = [
        {page_builder_entries}
    ]
    for page_no, builder in page_builders:
        for layer_spec in SLIDE_LAYER_SPECS:
            slide = prs.slides.add_slide(blank_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
            add_page_content(slide, page_no, builder, layer_spec)
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


if __name__ == "__main__":
    path = build_deck()
    print(path.resolve())
"""

def normalize_asset_adjustments(adjustments: Any) -> dict[str, Any]:
    if not isinstance(adjustments, dict):
        return {}

    normalized: dict[str, Any] = {}
    global_adjustment = _normalize_numeric_mapping(adjustments.get("global"), ASSET_RELATIVE_FIELDS)
    if global_adjustment:
        normalized["global"] = global_adjustment

    asset_map: dict[str, dict[str, int]] = {}
    raw_asset_map = adjustments.get("asset_map")
    if isinstance(raw_asset_map, dict):
        for raw_index, raw_plan in raw_asset_map.items():
            index = _coerce_positive_int(raw_index)
            if index is None:
                continue
            plan = _normalize_single_asset_adjustment(raw_plan)
            if plan:
                asset_map[str(index)] = plan

    raw_assets = adjustments.get("assets")
    if isinstance(raw_assets, list):
        for item in raw_assets:
            if not isinstance(item, dict):
                continue
            index = _coerce_positive_int(item.get("index"))
            if index is None:
                continue
            plan = _normalize_single_asset_adjustment(item)
            if plan:
                asset_map[str(index)] = plan

    if asset_map:
        normalized["asset_map"] = asset_map
    return normalized


def _build_page_function_source(page_no: int, script: str) -> str:
    body = script or "pass"
    page_texts_line = f'page_texts = PAGE_TEXTS["{page_no}"]'
    full_body = "\n".join([page_texts_line, body]) if body != "pass" else f"{page_texts_line}\npass"
    indented = "\n".join(f"    {line}" if line else "" for line in full_body.splitlines()) or "    pass"
    return f"def build_slide_{page_no:02d}(slide):\n{indented}"


def _build_page_text_inventory(page: dict[str, Any]) -> dict[str, str]:
    inventory: dict[str, str] = {}
    title = str(page.get("title", "")).strip()
    if title:
        inventory["title"] = title

    bullets = page.get("bullets", [])
    if isinstance(bullets, list):
        normalized_bullets = [str(item).strip() for item in bullets if str(item).strip()]
        for index, bullet in enumerate(normalized_bullets, start=1):
            inventory[f"bullet_{index}"] = bullet

    if not inventory.get("bullet_1"):
        body_index = 1
        for item in page.get("texts", []):
            if not isinstance(item, dict):
                continue
            text = str(item.get("text", "")).strip()
            role = str(item.get("role", "")).strip().lower()
            if not text or role == "title":
                continue
            inventory[f"body_{body_index}"] = text
            body_index += 1

    summary = str(page.get("summary", "")).strip()
    if summary:
        inventory["summary"] = summary

    for item in page.get("texts", []):
        if not isinstance(item, dict):
            continue
        text = str(item.get("text", "")).strip()
        if not text:
            continue
        role = str(item.get("role", "")).strip().lower() or "text"
        key_prefix = role if role not in inventory else f"{role}_text"
        next_index = 1
        candidate_key = f"{key_prefix}_{next_index}"
        while candidate_key in inventory:
            next_index += 1
            candidate_key = f"{key_prefix}_{next_index}"
        inventory[candidate_key] = text

    if not inventory:
        inventory["summary"] = summary
    return inventory


def _normalize_single_asset_adjustment(raw_plan: Any) -> dict[str, int]:
    if not isinstance(raw_plan, dict):
        return {}

    normalized = _normalize_numeric_mapping(raw_plan, ASSET_RELATIVE_FIELDS)
    for source_name, target_name in ASSET_ABSOLUTE_FIELDS:
        if target_name in normalized:
            continue
        value = _coerce_int(raw_plan.get(source_name))
        if value is None:
            continue
        if target_name in {"width", "height"} and value <= 0:
            continue
        normalized[target_name] = value
    return normalized


def _normalize_numeric_mapping(raw_mapping: Any, fields: tuple[str, ...]) -> dict[str, int]:
    if not isinstance(raw_mapping, dict):
        return {}

    normalized: dict[str, int] = {}
    for field in fields:
        value = _coerce_int(raw_mapping.get(field))
        if value is None or value == 0:
            continue
        normalized[field] = value
    return normalized


def _coerce_positive_int(value: Any) -> int | None:
    result = _coerce_int(value)
    if result is None or result <= 0:
        return None
    return result


def _coerce_int(value: Any) -> int | None:
    try:
        return int(round(float(value)))
    except (TypeError, ValueError):
        return None
