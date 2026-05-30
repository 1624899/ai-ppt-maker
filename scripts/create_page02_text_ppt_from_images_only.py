from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


SLIDE_WIDTH_INCH = 13.333333
SLIDE_HEIGHT_INCH = 7.5
IMAGE_WIDTH_PX = 2048
IMAGE_HEIGHT_PX = 1152


@dataclass(frozen=True)
class TextBoxSpec:
    text: str
    left: int
    top: int
    width: int
    height: int
    font_size: int
    color: str
    bold: bool = False
    align: str = "LEFT"
    valign: str = "TOP"
    font_name: str = "Microsoft YaHei"


# 仅依据原稿图和元素图人工还原的文字层信息。
PAGE02_TEXT_BOXES = [
    TextBoxSpec("跳出聊天框", 72, 42, 620, 126, 58, "102A72", True),
    TextBoxSpec("01", 104, 306, 98, 88, 38, "FFFFFF", True, "CENTER", "MIDDLE"),
    TextBoxSpec("传统画图", 236, 322, 292, 66, 30, "122B75", True),
    TextBoxSpec("•", 110, 426, 26, 34, 24, "1457DA", True),
    TextBoxSpec("Draw.io 手动拉框", 144, 426, 320, 40, 20, "1C2D63"),
    TextBoxSpec("•", 110, 476, 26, 34, 24, "1457DA", True),
    TextBoxSpec("对齐连线耗时", 144, 476, 260, 40, 20, "1C2D63"),
    TextBoxSpec("•", 110, 526, 26, 34, 24, "1457DA", True),
    TextBoxSpec("流程维护成本高", 144, 526, 300, 40, 20, "1C2D63"),
    TextBoxSpec("效率风险", 154, 804, 158, 34, 18, "F0463E", True, "CENTER", "MIDDLE"),
    TextBoxSpec("02", 786, 306, 98, 88, 38, "FFFFFF", True, "CENTER", "MIDDLE"),
    TextBoxSpec("AI 转换", 918, 322, 246, 66, 30, "122B75", True),
    TextBoxSpec("•", 792, 426, 26, 34, 24, "1457DA", True),
    TextBoxSpec("输入业务逻辑文字", 826, 426, 320, 40, 20, "1C2D63"),
    TextBoxSpec("•", 792, 476, 26, 34, 24, "1457DA", True),
    TextBoxSpec("指令：转化为 Draw.io XML", 826, 476, 410, 40, 20, "1C2D63"),
    TextBoxSpec("•", 792, 526, 26, 34, 24, "1457DA", True),
    TextBoxSpec("AI 生成结构化代码", 826, 526, 340, 40, 20, "1C2D63"),
    TextBoxSpec("03", 1428, 306, 98, 88, 38, "FFFFFF", True, "CENTER", "MIDDLE"),
    TextBoxSpec("自动生成", 1554, 322, 270, 66, 30, "122B75", True),
    TextBoxSpec("•", 1432, 426, 26, 34, 24, "1457DA", True),
    TextBoxSpec("复制 XML 到 Draw.io", 1468, 426, 360, 40, 20, "1C2D63"),
    TextBoxSpec("•", 1432, 476, 26, 34, 24, "1457DA", True),
    TextBoxSpec("一秒生成整齐流程图", 1468, 476, 320, 40, 20, "1C2D63"),
    TextBoxSpec("•", 1432, 526, 26, 34, 24, "1457DA", True),
    TextBoxSpec("结构化思维具象化", 1468, 526, 320, 40, 20, "1C2D63"),
    TextBoxSpec("进阶玩法：让 AI 从聊天助手变成流程生产工具", 462, 982, 1330, 92, 34, "165AEC", True, "LEFT", "MIDDLE"),
]


def pixels_to_slide_x(value: int) -> float:
    return value / IMAGE_WIDTH_PX * SLIDE_WIDTH_INCH


def pixels_to_slide_y(value: int) -> float:
    return value / IMAGE_HEIGHT_PX * SLIDE_HEIGHT_INCH


def pixels_to_slide_w(value: int) -> float:
    return value / IMAGE_WIDTH_PX * SLIDE_WIDTH_INCH


def pixels_to_slide_h(value: int) -> float:
    return value / IMAGE_HEIGHT_PX * SLIDE_HEIGHT_INCH


def parse_rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color.strip().lstrip("#").upper())


def parse_align(value: str) -> PP_ALIGN:
    return getattr(PP_ALIGN, value.upper(), PP_ALIGN.LEFT)


def parse_anchor(value: str) -> MSO_ANCHOR:
    return getattr(MSO_ANCHOR, value.upper(), MSO_ANCHOR.TOP)


def add_background_image(slide, image_path: Path) -> None:
    slide.shapes.add_picture(
        str(image_path),
        Inches(0),
        Inches(0),
        width=Inches(SLIDE_WIDTH_INCH),
        height=Inches(SLIDE_HEIGHT_INCH),
    )


def apply_font(run, spec: TextBoxSpec) -> None:
    run.font.name = spec.font_name
    run.font.size = Pt(spec.font_size)
    run.font.bold = spec.bold
    run.font.color.rgb = parse_rgb(spec.color)

    # 明确设置中西文字体，减少 Office 打开后的字体回退。
    r_pr = run._r.get_or_add_rPr()
    r_pr.set(qn("a:ea"), spec.font_name)
    r_pr.set(qn("a:cs"), spec.font_name)
    latin = r_pr.get_or_add_latin()
    latin.typeface = spec.font_name


def add_text_box(slide, spec: TextBoxSpec) -> None:
    shape = slide.shapes.add_textbox(
        Inches(pixels_to_slide_x(spec.left)),
        Inches(pixels_to_slide_y(spec.top)),
        Inches(pixels_to_slide_w(spec.width)),
        Inches(pixels_to_slide_h(spec.height)),
    )

    # 文字框保持透明，避免覆盖元素图。
    shape.fill.background()
    shape.line.fill.background()

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = False
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    text_frame.vertical_anchor = parse_anchor(spec.valign)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = parse_align(spec.align)
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0

    run = paragraph.add_run()
    run.text = spec.text
    apply_font(run, spec)


def build_presentation(reference_image: Path, elements_image: Path, output_pptx: Path) -> Path:
    if not reference_image.exists():
        raise FileNotFoundError(f"原稿图不存在: {reference_image}")
    if not elements_image.exists():
        raise FileNotFoundError(f"元素图不存在: {elements_image}")

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_INCH)
    presentation.slide_height = Inches(SLIDE_HEIGHT_INCH)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    add_background_image(slide, elements_image)

    for spec in PAGE02_TEXT_BOXES:
        add_text_box(slide, spec)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_pptx)
    return output_pptx


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="仅基于两张输入图片重建 page 02 的文字 PPT。")
    parser.add_argument(
        "--reference-image",
        default=r"output\758b2965960a\01_reference_pages\page_02_reference.png",
        help="完整原稿图路径，仅用于核对输入来源。",
    )
    parser.add_argument(
        "--elements-image",
        default=r"output\758b2965960a\02_elements_pages\page_02_elements.png",
        help="元素图路径。",
    )
    parser.add_argument(
        "--output-pptx",
        default=r"output\758b2965960a\page_02_text_overlay_from_images_only.pptx",
        help="输出 PPTX 路径。",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output_path = build_presentation(
        reference_image=Path(args.reference_image),
        elements_image=Path(args.elements_image),
        output_pptx=Path(args.output_pptx),
    )
    print(output_path.resolve())


if __name__ == "__main__":
    main()
