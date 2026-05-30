from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


SLIDE_WIDTH_INCH = 13.333333
SLIDE_HEIGHT_INCH = 7.5
IMAGE_WIDTH_PX = 2048
IMAGE_HEIGHT_PX = 1152


@dataclass(frozen=True)
class TextSpec:
    text: str
    left: int
    top: int
    width: int
    height: int
    font_size: int
    color: str
    bold: bool = False
    align: str = "LEFT"
    valign: str = "TOP"
    font_name: str = "Microsoft YaHei"


PAGE02_TEXT_SPECS = [
    TextSpec("跳出聊天框", 72, 18, 720, 132, 60, "0A2369", True),
    TextSpec("01", 102, 304, 100, 90, 38, "FFFFFF", True, "CENTER", "MIDDLE"),
    TextSpec("传统画图", 238, 320, 260, 68, 30, "163A63", True),
    TextSpec("• Draw.io 手动拉框", 112, 430, 360, 42, 20, "1D2C63"),
    TextSpec("• 对齐连线耗时", 112, 479, 360, 42, 20, "1D2C63"),
    TextSpec("• 流程维护成本高", 112, 528, 390, 42, 20, "1D2C63"),
    TextSpec("效率风险", 152, 804, 160, 34, 18, "E54444", True, "CENTER", "MIDDLE"),
    TextSpec("02", 784, 304, 100, 90, 38, "FFFFFF", True, "CENTER", "MIDDLE"),
    TextSpec("AI 转换", 916, 320, 250, 68, 30, "163A63", True),
    TextSpec("• 输入业务逻辑文字", 825, 430, 360, 42, 20, "1D2C63"),
    TextSpec("• 指令：转化为 Draw.io XML", 825, 479, 430, 42, 20, "1D2C63"),
    TextSpec("• AI 生成结构化代码", 825, 528, 390, 42, 20, "1D2C63"),
    TextSpec("03", 1425, 304, 100, 90, 38, "FFFFFF", True, "CENTER", "MIDDLE"),
    TextSpec("自动生成", 1550, 320, 260, 68, 30, "163A63", True),
    TextSpec("• 复制 XML 到 Draw.io", 1464, 430, 400, 42, 20, "1D2C63"),
    TextSpec("• 一秒生成整齐流程图", 1464, 479, 400, 42, 20, "1D2C63"),
    TextSpec("• 结构化思维具象化", 1464, 528, 400, 42, 20, "1D2C63"),
    TextSpec("进阶玩法：让 AI 从聊天助手变成流程生产工具", 462, 982, 1360, 92, 34, "175BEA", True, "LEFT", "MIDDLE"),
]


def px_to_inches_x(value: int) -> float:
    return value / IMAGE_WIDTH_PX * SLIDE_WIDTH_INCH


def px_to_inches_y(value: int) -> float:
    return value / IMAGE_HEIGHT_PX * SLIDE_HEIGHT_INCH


def px_to_inches_w(value: int) -> float:
    return value / IMAGE_WIDTH_PX * SLIDE_WIDTH_INCH


def px_to_inches_h(value: int) -> float:
    return value / IMAGE_HEIGHT_PX * SLIDE_HEIGHT_INCH


def to_rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color.strip().lstrip("#").upper())


def to_align(value: str) -> PP_ALIGN:
    return getattr(PP_ALIGN, value.upper(), PP_ALIGN.LEFT)


def to_anchor(value: str) -> MSO_ANCHOR:
    return getattr(MSO_ANCHOR, value.upper(), MSO_ANCHOR.TOP)


def add_full_slide_image(slide, image_path: Path) -> None:
    slide.shapes.add_picture(
        str(image_path),
        Inches(0),
        Inches(0),
        width=Inches(SLIDE_WIDTH_INCH),
        height=Inches(SLIDE_HEIGHT_INCH),
    )


def add_textbox(slide, spec: TextSpec) -> None:
    shape = slide.shapes.add_textbox(
        Inches(px_to_inches_x(spec.left)),
        Inches(px_to_inches_y(spec.top)),
        Inches(px_to_inches_w(spec.width)),
        Inches(px_to_inches_h(spec.height)),
    )
    shape.fill.background()
    shape.line.fill.background()

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = False
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    text_frame.vertical_anchor = to_anchor(spec.valign)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = to_align(spec.align)
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.0

    run = paragraph.add_run()
    run.text = spec.text
    run.font.name = spec.font_name
    run.font.size = Pt(spec.font_size)
    run.font.bold = spec.bold
    run.font.color.rgb = to_rgb(spec.color)

    # 补齐东亚字体设置，避免 PowerPoint 回退到其它字体。
    r_pr = run._r.get_or_add_rPr()
    r_pr.set(qn("a:ea"), spec.font_name)
    r_pr.set(qn("a:cs"), spec.font_name)
    latin = r_pr.get_or_add_latin()
    latin.typeface = spec.font_name


def build_ppt(reference_image: Path, elements_image: Path, output_pptx: Path) -> Path:
    if not reference_image.exists():
        raise FileNotFoundError(f"原稿图不存在: {reference_image}")
    if not elements_image.exists():
        raise FileNotFoundError(f"元素图不存在: {elements_image}")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCH)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_full_slide_image(slide, elements_image)
    for spec in PAGE02_TEXT_SPECS:
        add_textbox(slide, spec)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return output_pptx


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="生成仅包含元素图和文字层的 page 02 PPT。")
    parser.add_argument(
        "--reference-image",
        default=r"output\758b2965960a\01_reference_pages\page_02_reference.png",
        help="参考页图片路径，仅用于校验输入完整性。",
    )
    parser.add_argument(
        "--elements-image",
        default=r"output\758b2965960a\02_elements_pages\page_02_elements.png",
        help="元素页图片路径，将整页铺到底图。",
    )
    parser.add_argument(
        "--output-pptx",
        default=r"output\758b2965960a\page_02_text_overlay.pptx",
        help="输出 PPTX 路径。",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output_path = build_ppt(
        reference_image=Path(args.reference_image),
        elements_image=Path(args.elements_image),
        output_pptx=Path(args.output_pptx),
    )
    print(output_path.resolve())


if __name__ == "__main__":
    main()
