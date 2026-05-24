from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
REFERENCE_IMAGE = ROOT / "output" / "758b2965960a" / "01_reference_pages" / "page_02_reference.png"
ELEMENTS_IMAGE = ROOT / "output" / "758b2965960a" / "02_elements_pages" / "page_02_elements.png"
OUTPUT_PPTX = ROOT / "output" / "758b2965960a" / "page_02_text_only_merged.pptx"

PX_WIDTH = 2048
PX_HEIGHT = 1152
SLIDE_WIDTH_IN = 16
SLIDE_HEIGHT_IN = 9

BLUE = RGBColor(0, 20, 82)
BRIGHT_BLUE = RGBColor(0, 82, 214)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(255, 32, 32)


def px_to_in(value: float, axis: str) -> float:
    """把参考图像素坐标换算成幻灯片英寸坐标。"""
    base_px = PX_WIDTH if axis == "x" else PX_HEIGHT
    base_in = SLIDE_WIDTH_IN if axis == "x" else SLIDE_HEIGHT_IN
    return value / base_px * base_in


def add_textbox(slide, text, x, y, w, h, font_size, color=BLUE, bold=False, align=PP_ALIGN.LEFT):
    """按像素坐标添加无填充、无线条文本框。"""
    shape = slide.shapes.add_textbox(
        Inches(px_to_in(x, "x")),
        Inches(px_to_in(y, "y")),
        Inches(px_to_in(w, "x")),
        Inches(px_to_in(h, "y")),
    )
    shape.fill.background()
    shape.line.fill.background()
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_bullets(slide, items, x, y, w, line_gap, font_size=19):
    """添加项目符号和文字，保持每一行独立便于微调。"""
    for index, item in enumerate(items):
        top = y + index * line_gap
        add_textbox(slide, "•", x, top - 2, 22, 34, 26, color=BRIGHT_BLUE, bold=True)
        add_textbox(slide, item, x + 38, top, w, 34, font_size, color=BLUE)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(ELEMENTS_IMAGE), 0, 0, width=prs.slide_width, height=prs.slide_height)

    add_textbox(slide, "跳出聊天框", 72, 42, 620, 125, 50, color=BLUE, bold=True)

    card_data = [
        {
            "num": "01",
            "num_box": (118, 319, 100, 70),
            "title": "传统画图",
            "title_box": (237, 325, 280, 60),
            "bullets": ["Draw.io 手动拉框", "对齐连线耗时", "流程维护成本高"],
            "bullet_origin": (113, 429),
        },
        {
            "num": "02",
            "num_box": (799, 319, 100, 70),
            "title": "AI 转换",
            "title_box": (916, 326, 300, 60),
            "bullets": ["输入业务逻辑文字", "指令：转化为 Draw.io XML", "AI 生成结构化代码"],
            "bullet_origin": (792, 429),
        },
        {
            "num": "03",
            "num_box": (1447, 319, 100, 70),
            "title": "自动生成",
            "title_box": (1555, 326, 310, 60),
            "bullets": ["复制 XML 到 Draw.io", "一秒生成整齐流程图", "结构化思维具象化"],
            "bullet_origin": (1438, 429),
        },
    ]

    for card in card_data:
        add_textbox(slide, card["num"], *card["num_box"], 31, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, card["title"], *card["title_box"], 28, color=BLUE, bold=True)
        add_bullets(slide, card["bullets"], card["bullet_origin"][0], card["bullet_origin"][1], 520, 49)

    add_textbox(slide, "效率风险", 174, 805, 128, 42, 19, color=RED, bold=True)
    add_textbox(slide, "AI", 1138, 685, 88, 72, 38, color=BRIGHT_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, "进阶玩法：让 AI 从聊天助手变成流程生产工具", 466, 987, 1180, 72, 31, color=BRIGHT_BLUE, bold=True)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


if __name__ == "__main__":
    if not REFERENCE_IMAGE.exists():
        raise FileNotFoundError(f"缺少参考图：{REFERENCE_IMAGE}")
    if not ELEMENTS_IMAGE.exists():
        raise FileNotFoundError(f"缺少元素图：{ELEMENTS_IMAGE}")
    print(build_presentation())
